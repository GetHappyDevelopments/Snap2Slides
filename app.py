from __future__ import annotations

import argparse
import ctypes
import json
import os
import sys
import tempfile
import threading
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

import cv2
import numpy as np
from lxml import etree
from PIL import Image, ImageDraw, ImageTk
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Pt
from rapidocr_onnxruntime import RapidOCR
import tkinter as tk
from tkinter import filedialog, messagebox, ttk


EMU_PER_INCH = 914400
SLIDE_WIDTH_IN = 13.333333
SLIDE_HEIGHT_IN = 7.5
MIN_TEXT_CONFIDENCE = 0.45
OCR_IMAGE_SCALE = 3
EXPORT_MODES = ("pixel", "hybrid", "vector", "visual_safe", "editable", "styled_reconstruct", "reconstruct")
BURGUNDY = RGBColor(137, 13, 64)
REFERENCE_ACCENT = RGBColor(154, 14, 62)
REFERENCE_TEXT = RGBColor(17, 17, 17)
REFERENCE_LOGO = RGBColor(21, 93, 109)
TEXT_BLACK = RGBColor(20, 20, 20)
APP_ROOT = Path(__file__).resolve().parent
APP_ICON_ICO = APP_ROOT / "assets" / "app_icon.ico"
APP_ICON_PNG = APP_ROOT / "assets" / "app_icon.png"
APP_SPLASH_PNG = APP_ROOT / "assets" / "splash.png"
WINDOWS_APP_ID = "Snap2Slides.Image2PptSlicer"


@dataclass(frozen=True)
class Rect:
    x: int
    y: int
    w: int
    h: int

    @property
    def x2(self) -> int:
        return self.x + self.w

    @property
    def y2(self) -> int:
        return self.y + self.h

    @property
    def area(self) -> int:
        return self.w * self.h

    def clipped(self, max_w: int, max_h: int) -> "Rect":
        x = max(0, min(self.x, max_w - 1))
        y = max(0, min(self.y, max_h - 1))
        x2 = max(x + 1, min(self.x2, max_w))
        y2 = max(y + 1, min(self.y2, max_h))
        return Rect(x, y, x2 - x, y2 - y)

    def padded(self, px: int, max_w: int, max_h: int) -> "Rect":
        return Rect(self.x - px, self.y - px, self.w + px * 2, self.h + px * 2).clipped(max_w, max_h)

    def moved(self, x: int, y: int, max_w: int, max_h: int) -> "Rect":
        x = max(0, min(int(x), max_w - self.w))
        y = max(0, min(int(y), max_h - self.h))
        return Rect(x, y, self.w, self.h)

    def resized(self, w: int, h: int, max_w: int, max_h: int) -> "Rect":
        w = max(6, min(int(w), max_w - self.x))
        h = max(6, min(int(h), max_h - self.y))
        return Rect(self.x, self.y, w, h)


@dataclass(frozen=True)
class ReconstructionConfig:
    output_slide_format: str = "16:9"
    slide_width_inches: float = SLIDE_WIDTH_IN
    slide_height_inches: float = SLIDE_HEIGHT_IN
    border_trim_px: int = 8
    border_trim_ratio: float = 0.012
    min_slide_area: int = 9000
    max_slide_area_ratio: float = 0.98
    expected_aspect_ratios: tuple[float, ...] = (16 / 9, 4 / 3)
    aspect_ratio_tolerance: float = 0.32
    vector_score_threshold: float = 0.56
    raster_score_threshold: float = 0.58
    ocr_confidence_threshold: float = MIN_TEXT_CONFIDENCE
    enable_svg_export: bool = True
    enable_native_shapes: bool = True
    enable_raster_fallback: bool = True
    enable_debug_output: bool = False
    enable_perspective_correction: bool = True
    enable_com_powerpoint_export: bool = False
    fallback_to_pixel_mode: bool = True


DEFAULT_CONFIG = ReconstructionConfig()


@dataclass(frozen=True)
class SlideCrop:
    outer_rect: Rect
    inner_rect: Rect


@dataclass
class TextLine:
    rect: Rect
    text: str
    confidence: float
    color: RGBColor


@dataclass
class TextBlock:
    rect: Rect
    lines: list[TextLine]
    bullet_lines: list[bool]
    order: int = 0

    @property
    def text(self) -> str:
        return "\n".join(line.text for line in self.lines)

    @property
    def confidence(self) -> float:
        return min(line.confidence for line in self.lines)

    @property
    def color(self) -> RGBColor:
        return self.lines[0].color if self.lines else TEXT_BLACK


@dataclass
class VisualBlock:
    rect: Rect
    image: Image.Image
    kind: str = "visual"


@dataclass
class ShapeBlock:
    rect: Rect
    fill: RGBColor


@dataclass
class VectorElement:
    kind: str
    rect: Rect
    score: float
    fill: RGBColor | None = None
    stroke: RGBColor | None = None
    stroke_width: int = 1
    points: list[tuple[int, int]] | None = None
    reasons: list[str] | None = None


@dataclass
class AnalyzedSlide:
    rect: Rect
    image: Image.Image
    texts: list[TextBlock]
    visuals: list[VisualBlock]
    shapes: list[ShapeBlock]
    vectors: list[VectorElement] | None = None
    background: RGBColor | None = None
    source_rect: Rect | None = None
    debug: dict | None = None


@dataclass(frozen=True)
class TextStyle:
    size: int
    color: RGBColor
    bold: bool = False
    font_name: str = "Aptos"


@dataclass(frozen=True)
class StyleProfile:
    cover_title: TextStyle
    cover_subtitle: TextStyle
    title: TextStyle
    heading: TextStyle
    card_heading: TextStyle
    body: TextStyle
    compact_body: TextStyle
    bullet: TextStyle
    bullet_marker: TextStyle
    logo: TextStyle


REFERENCE_STYLE = StyleProfile(
    cover_title=TextStyle(18, REFERENCE_ACCENT, bold=True, font_name="Aptos Display"),
    cover_subtitle=TextStyle(10, REFERENCE_TEXT),
    title=TextStyle(21, REFERENCE_TEXT, bold=True, font_name="Aptos Display"),
    heading=TextStyle(21, REFERENCE_TEXT, bold=True, font_name="Aptos Display"),
    card_heading=TextStyle(8, REFERENCE_TEXT, bold=True),
    body=TextStyle(12, REFERENCE_TEXT),
    compact_body=TextStyle(7, REFERENCE_TEXT),
    bullet=TextStyle(12, REFERENCE_TEXT),
    bullet_marker=TextStyle(14, REFERENCE_ACCENT, bold=True),
    logo=TextStyle(21, REFERENCE_LOGO, bold=True, font_name="Aptos Display"),
)


def _rgb_to_hex(color: RGBColor | None, fallback: str = "none") -> str:
    if color is None:
        return fallback
    return f"#{int(color[0]):02x}{int(color[1]):02x}{int(color[2]):02x}"


def _rgb_to_tuple(color: RGBColor | None) -> tuple[int, int, int] | None:
    if color is None:
        return None
    return int(color[0]), int(color[1]), int(color[2])


def load_config(path: Path | None) -> ReconstructionConfig:
    if path is None:
        return DEFAULT_CONFIG
    data = json.loads(path.read_text(encoding="utf-8"))
    allowed = ReconstructionConfig.__dataclass_fields__
    values = {key: value for key, value in data.items() if key in allowed}
    if "expected_aspect_ratios" in values:
        values["expected_aspect_ratios"] = tuple(float(item) for item in values["expected_aspect_ratios"])
    return ReconstructionConfig(**values)


def _bands_from_mask(values: np.ndarray, min_width: int) -> list[tuple[int, int]]:
    bands: list[tuple[int, int]] = []
    start: int | None = None
    for idx, is_band in enumerate(values):
        if is_band and start is None:
            start = idx
        elif not is_band and start is not None:
            if idx - start >= min_width:
                bands.append((start, idx))
            start = None
    if start is not None and len(values) - start >= min_width:
        bands.append((start, len(values)))
    return bands


def _intervals_between_bands(size: int, bands: list[tuple[int, int]], min_size: int) -> list[tuple[int, int]]:
    intervals: list[tuple[int, int]] = []
    cursor = 0
    for start, end in bands:
        if start - cursor >= min_size:
            intervals.append((cursor, start))
        cursor = end
    if size - cursor >= min_size:
        intervals.append((cursor, size))
    return intervals


def _smoothed_axis_bands(gray: np.ndarray, axis: int, threshold: float, min_width: int) -> list[tuple[int, int]]:
    values = gray.mean(axis=axis)
    kernel_width = max(9, min_width * 3)
    kernel = np.ones(kernel_width) / kernel_width
    smoothed = np.convolve(values, kernel, mode="same")
    return _bands_from_mask(smoothed < threshold, min_width)


def _is_expected_slide_aspect(rect: Rect, config: ReconstructionConfig) -> bool:
    aspect = rect.w / max(1, rect.h)
    return any(abs(aspect - expected) <= config.aspect_ratio_tolerance for expected in config.expected_aspect_ratios)


def _sort_slide_rects(rects: list[Rect]) -> list[Rect]:
    if not rects:
        return []
    median_h = float(np.median([rect.h for rect in rects]))
    row_tolerance = max(12, int(median_h * 0.35))
    rows: list[list[Rect]] = []
    for rect in sorted(rects, key=lambda item: item.y):
        for row in rows:
            if abs(rect.y - row[0].y) <= row_tolerance:
                row.append(rect)
                break
        else:
            rows.append([rect])
    sorted_rects: list[Rect] = []
    for row in rows:
        sorted_rects.extend(sorted(row, key=lambda item: item.x))
    return sorted_rects


def _dedupe_rect_candidates(rects: list[Rect]) -> list[Rect]:
    kept: list[Rect] = []
    for candidate in sorted(rects, key=lambda item: item.area, reverse=True):
        if any(_overlap_ratio(candidate, existing) > 0.78 for existing in kept):
            continue
        kept.append(candidate)
    return _sort_slide_rects(kept)


def _detect_slide_rects_by_contours(gray: np.ndarray, config: ReconstructionConfig) -> list[Rect]:
    h, w = gray.shape
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blurred, 50, 150)
    edges = cv2.dilate(edges, np.ones((3, 3), np.uint8), iterations=1)
    contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    candidates: list[Rect] = []
    max_area = w * h * config.max_slide_area_ratio
    for contour in contours:
        peri = cv2.arcLength(contour, True)
        approx = cv2.approxPolyDP(contour, 0.025 * peri, True)
        x, y, bw, bh = cv2.boundingRect(approx)
        rect = Rect(int(x), int(y), int(bw), int(bh)).clipped(w, h)
        if rect.area < config.min_slide_area or rect.area > max_area:
            continue
        if rect.w < 80 or rect.h < 60 or not _is_expected_slide_aspect(rect, config):
            continue
        extent = cv2.contourArea(contour) / max(1, rect.area)
        if len(approx) >= 4 and extent > 0.55:
            candidates.append(rect)
    return _dedupe_rect_candidates(candidates)


def detect_slide_rects(image: Image.Image, config: ReconstructionConfig | None = None) -> list[Rect]:
    config = config or DEFAULT_CONFIG
    arr = np.array(image.convert("RGB"))
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    h, w = gray.shape

    contour_rects = _detect_slide_rects_by_contours(gray, config)
    if len(contour_rects) >= 1:
        return contour_rects

    min_band_w = max(6, w // 300)
    min_band_h = max(6, h // 300)
    col_bands = _smoothed_axis_bands(gray, axis=0, threshold=132, min_width=min_band_w)
    row_bands = _smoothed_axis_bands(gray, axis=1, threshold=132, min_width=min_band_h)

    if len(col_bands) < 2 or len(row_bands) < 2:
        dark = gray < 95
        vertical = dark.mean(axis=0) > 0.25
        horizontal = dark.mean(axis=1) > 0.25
        col_bands = _bands_from_mask(vertical, min_band_w)
        row_bands = _bands_from_mask(horizontal, min_band_h)

    min_slide_w = max(120, w // 8)
    min_slide_h = max(80, h // 8)
    col_intervals = _intervals_between_bands(w, col_bands, min_slide_w)
    row_intervals = _intervals_between_bands(h, row_bands, min_slide_h)

    rects = [
        Rect(x1, y1, x2 - x1, y2 - y1).padded(-2, w, h)
        for y1, y2 in row_intervals
        for x1, x2 in col_intervals
        if (x2 - x1) > min_slide_w and (y2 - y1) > min_slide_h
    ]

    rects = [rect for rect in rects if _is_expected_slide_aspect(rect, config)]
    if len(rects) >= 2:
        return _sort_slide_rects(rects)

    return [Rect(0, 0, w, h)]


def _box_to_rect(points: list[list[float]], max_w: int, max_h: int) -> Rect:
    xs = [int(round(p[0])) for p in points]
    ys = [int(round(p[1])) for p in points]
    return Rect(min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys)).padded(2, max_w, max_h)


def _union_rect(rects: list[Rect], max_w: int, max_h: int, padding: int = 0) -> Rect:
    x1 = min(rect.x for rect in rects)
    y1 = min(rect.y for rect in rects)
    x2 = max(rect.x2 for rect in rects)
    y2 = max(rect.y2 for rect in rects)
    return Rect(x1, y1, x2 - x1, y2 - y1).padded(padding, max_w, max_h)


def _inner_slide_crop(rect: Rect, source_w: int, source_h: int, config: ReconstructionConfig) -> SlideCrop:
    trim = max(config.border_trim_px, int(min(rect.w, rect.h) * config.border_trim_ratio))
    inner = Rect(rect.x + trim, rect.y + trim, rect.w - trim * 2, rect.h - trim * 2).clipped(source_w, source_h)
    if inner.w < 10 or inner.h < 10:
        inner = rect.clipped(source_w, source_h)
    return SlideCrop(outer_rect=rect.clipped(source_w, source_h), inner_rect=inner)


def _guess_text_color(arr: np.ndarray, rect: Rect) -> RGBColor:
    crop = arr[rect.y : rect.y2, rect.x : rect.x2]
    if crop.size == 0:
        return TEXT_BLACK
    darkish = crop[np.mean(crop, axis=2) < 150]
    if darkish.size == 0:
        return TEXT_BLACK
    mean = darkish.mean(axis=0)
    if mean[0] > mean[1] * 1.25 and mean[0] > mean[2] * 1.15:
        return BURGUNDY
    return TEXT_BLACK


def _intersection_area(first: Rect, second: Rect) -> int:
    x1 = max(first.x, second.x)
    y1 = max(first.y, second.y)
    x2 = min(first.x2, second.x2)
    y2 = min(first.y2, second.y2)
    if x2 <= x1 or y2 <= y1:
        return 0
    return (x2 - x1) * (y2 - y1)


def _overlap_ratio(first: Rect, second: Rect) -> float:
    smaller = min(first.area, second.area)
    if smaller <= 0:
        return 0.0
    return _intersection_area(first, second) / smaller


def _merge_rects(rects: list[Rect], overlap_threshold: float = 0.2) -> list[Rect]:
    if not rects:
        return []

    merged: list[Rect] = []
    for rect in sorted(rects, key=lambda item: (item.y, item.x)):
        for index, existing in enumerate(merged):
            intersection = _intersection_area(rect, existing)
            union = rect.area + existing.area - intersection
            if union > 0 and intersection / union > overlap_threshold:
                x1 = min(rect.x, existing.x)
                y1 = min(rect.y, existing.y)
                x2 = max(rect.x2, existing.x2)
                y2 = max(rect.y2, existing.y2)
                merged[index] = Rect(x1, y1, x2 - x1, y2 - y1)
                break
        else:
            merged.append(rect)
    return merged


def _rect_fill_color(arr: np.ndarray, rect: Rect) -> RGBColor:
    crop = arr[rect.y : rect.y2, rect.x : rect.x2]
    if crop.size == 0:
        return RGBColor(255, 255, 255)
    median = np.median(crop.reshape(-1, 3), axis=0)
    return RGBColor(int(median[0]), int(median[1]), int(median[2]))


def _build_text_mask(width: int, height: int, text_rects: list[Rect], padding: int = 2) -> np.ndarray:
    mask = np.zeros((height, width), dtype=np.uint8)
    for rect in text_rects:
        r = rect.padded(padding, width, height)
        mask[r.y : r.y2, r.x : r.x2] = 255
    return mask


def _build_text_pixel_mask(image: Image.Image, text_rects: list[Rect], padding: int = 6) -> np.ndarray:
    arr = np.array(image.convert("RGB"))
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(arr, cv2.COLOR_RGB2HSV)
    saturation = hsv[:, :, 1]
    mask = np.zeros((image.height, image.width), dtype=np.uint8)

    for rect in text_rects:
        r = rect.padded(padding, image.width, image.height)
        crop_gray = gray[r.y : r.y2, r.x : r.x2]
        crop_sat = saturation[r.y : r.y2, r.x : r.x2]
        if crop_gray.size == 0:
            continue

        median_gray = float(np.median(crop_gray))
        if median_gray < 120:
            local_pixels = (crop_gray > 155) | (np.abs(crop_gray.astype(np.int16) - int(median_gray)) > 55)
        else:
            local_pixels = (crop_gray < 185) | ((crop_gray < 225) & (crop_sat > 35))
        local = local_pixels.astype(np.uint8) * 255
        local = cv2.dilate(local, np.ones((3, 3), np.uint8), iterations=1)
        mask[r.y : r.y2, r.x : r.x2] = np.maximum(mask[r.y : r.y2, r.x : r.x2], local)

    return mask


def _remove_text_from_image(image: Image.Image, text_rects: list[Rect]) -> Image.Image:
    if not text_rects:
        return image

    arr = np.array(image.convert("RGB"))
    mask = _build_text_pixel_mask(image, text_rects, padding=6)
    if not np.any(mask):
        return image

    inpainted = cv2.inpaint(arr, mask, 3, cv2.INPAINT_TELEA)
    return Image.fromarray(inpainted)


def _component_rect_for_seed(mask: np.ndarray, seed_rect: Rect) -> Rect | None:
    num_labels, labels, stats, _centroids = cv2.connectedComponentsWithStats(mask, 8)
    if num_labels <= 1:
        return None

    y1 = max(0, seed_rect.y)
    y2 = min(mask.shape[0], seed_rect.y2)
    x1 = max(0, seed_rect.x)
    x2 = min(mask.shape[1], seed_rect.x2)
    seed_labels = labels[y1:y2, x1:x2]
    label_ids, counts = np.unique(seed_labels[seed_labels > 0], return_counts=True)
    if label_ids.size == 0:
        return None

    label = int(label_ids[int(np.argmax(counts))])
    x, y, w, h, _area = stats[label]
    return Rect(int(x), int(y), int(w), int(h))


def _detect_text_containers(slide_image: Image.Image, text_blocks: list[TextBlock], text_rects: list[Rect]) -> list[ShapeBlock]:
    if not text_blocks:
        return []

    cleaned_image = _remove_text_from_image(slide_image, text_rects)
    arr = np.array(cleaned_image.convert("RGB"))
    h, w = arr.shape[:2]
    shapes: list[ShapeBlock] = []

    for block in text_blocks:
        normalized_text = _normalized_text(block.text)
        if len(normalized_text) <= 5 or normalized_text in {"msg", "6sw", "msq"}:
            continue

        seed = block.rect.padded(max(4, block.rect.h // 4), w, h)
        seed_crop = arr[seed.y : seed.y2, seed.x : seed.x2]
        if seed_crop.size == 0:
            continue

        median = np.median(seed_crop.reshape(-1, 3), axis=0)
        seed_std = float(np.mean(np.std(seed_crop.reshape(-1, 3), axis=0)))
        if seed_std > 58:
            continue

        distance = np.linalg.norm(arr.astype(np.int16) - median.astype(np.int16), axis=2)
        threshold = 34 if seed_std < 22 else 46
        similar = (distance < threshold).astype(np.uint8) * 255
        similar = cv2.morphologyEx(similar, cv2.MORPH_CLOSE, np.ones((7, 7), np.uint8), iterations=2)
        similar = cv2.morphologyEx(similar, cv2.MORPH_OPEN, np.ones((3, 3), np.uint8), iterations=1)

        rect = _component_rect_for_seed(similar, seed)
        if rect is None:
            continue
        rect = rect.padded(1, w, h)
        if rect.area < block.rect.area * 1.08 or rect.area > w * h * 0.86:
            continue
        if rect.w < block.rect.w * 0.9 or rect.h < block.rect.h * 0.9:
            continue

        fill = _rect_fill_color(arr, rect)
        shapes.append(ShapeBlock(rect=rect, fill=fill))

    return _dedupe_shapes(shapes, w, h)


def _detect_visuals(slide_image: Image.Image, text_rects: list[Rect]) -> list[VisualBlock]:
    arr = np.array(slide_image.convert("RGB"))
    h, w = arr.shape[:2]
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(arr, cv2.COLOR_RGB2HSV)
    saturation = hsv[:, :, 1]

    text_mask = _build_text_mask(w, h, text_rects, padding=5)
    visual_source = _remove_text_from_image(slide_image, text_rects)

    # Keep non-white or colorful content, but remove OCR text so photos, logos, icons and colored panels remain.
    not_white = ((gray < 248) | (saturation > 28)).astype(np.uint8) * 255
    not_white[text_mask > 0] = 0

    kernel = np.ones((5, 5), np.uint8)
    cleaned = cv2.morphologyEx(not_white, cv2.MORPH_CLOSE, kernel, iterations=2)
    cleaned = cv2.morphologyEx(cleaned, cv2.MORPH_OPEN, np.ones((3, 3), np.uint8), iterations=1)

    num_labels, labels, stats, _ = cv2.connectedComponentsWithStats(cleaned, 8)
    min_area = max(700, int(w * h * 0.002))
    visuals: list[VisualBlock] = []

    for label in range(1, num_labels):
        x, y, bw, bh, area = stats[label]
        if area < min_area or bw < 18 or bh < 18:
            continue
        rect = Rect(int(x), int(y), int(bw), int(bh)).padded(3, w, h)
        crop = visual_source.crop((rect.x, rect.y, rect.x2, rect.y2))
        visuals.append(VisualBlock(rect=rect, image=crop, kind=_classify_visual(arr, rect)))

    return _dedupe_visuals(visuals, w, h)


def _is_mostly_white_region(arr: np.ndarray, rect: Rect) -> bool:
    crop = arr[rect.y : rect.y2, rect.x : rect.x2]
    if crop.size == 0:
        return True
    gray = cv2.cvtColor(crop, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(crop, cv2.COLOR_RGB2HSV)
    white_ratio = float(np.mean((gray > 244) & (hsv[:, :, 1] < 24)))
    return white_ratio > 0.94


def _is_large_visual_candidate(visual: VisualBlock, slide_w: int, slide_h: int) -> bool:
    area_ratio = visual.rect.area / max(1, slide_w * slide_h)
    return (
        visual.kind == "photo"
        or (area_ratio >= 0.055 and visual.rect.w >= slide_w * 0.18 and visual.rect.h >= slide_h * 0.16)
    )


def _detect_large_image_regions(
    slide_image: Image.Image,
    text_rects: list[Rect],
    visual_candidates: list[VisualBlock],
) -> list[VisualBlock]:
    arr = np.array(slide_image.convert("RGB"))
    h, w = arr.shape[:2]
    candidates = [
        VisualBlock(visual.rect, visual.image, "photo")
        for visual in visual_candidates
        if _is_large_visual_candidate(visual, w, h) and not _is_mostly_white_region(arr, visual.rect)
    ]
    for photo in _detect_layout_photo_regions(slide_image, text_rects):
        if any(_overlap_ratio(photo.rect, existing.rect) > 0.86 and photo.rect.area > existing.rect.area * 1.45 for existing in candidates):
            continue
        candidates.append(photo)
    return _dedupe_visuals(candidates, w, h)


def _alpha_crop_for_content(slide_image: Image.Image, rect: Rect, exclude_rects: list[Rect]) -> Image.Image:
    crop = np.array(slide_image.crop((rect.x, rect.y, rect.x2, rect.y2)).convert("RGB"))
    if crop.size == 0:
        return Image.fromarray(crop)
    gray = cv2.cvtColor(crop, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(crop, cv2.COLOR_RGB2HSV)
    alpha = ((gray < 244) | (hsv[:, :, 1] > 26)).astype(np.uint8) * 255
    for excluded in exclude_rects:
        ix1 = max(rect.x, excluded.x) - rect.x
        iy1 = max(rect.y, excluded.y) - rect.y
        ix2 = min(rect.x2, excluded.x2) - rect.x
        iy2 = min(rect.y2, excluded.y2) - rect.y
        if ix2 > ix1 and iy2 > iy1:
            alpha[iy1:iy2, ix1:ix2] = 0
    alpha = cv2.morphologyEx(alpha, cv2.MORPH_CLOSE, np.ones((2, 2), np.uint8), iterations=1)
    return Image.fromarray(np.dstack([crop, alpha]))


def _detect_icon_regions(slide_image: Image.Image, exclude_rects: list[Rect]) -> list[VisualBlock]:
    arr = np.array(slide_image.convert("RGB"))
    h, w = arr.shape[:2]
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(arr, cv2.COLOR_RGB2HSV)
    saturation = hsv[:, :, 1]
    mask = ((gray < 242) | (saturation > 30)).astype(np.uint8) * 255
    for rect in exclude_rects:
        padded = rect.padded(3, w, h)
        mask[padded.y : padded.y2, padded.x : padded.x2] = 0
    mask = cv2.morphologyEx(mask, cv2.MORPH_CLOSE, np.ones((4, 4), np.uint8), iterations=2)
    mask = cv2.morphologyEx(mask, cv2.MORPH_OPEN, np.ones((2, 2), np.uint8), iterations=1)

    num_labels, _labels, stats, _ = cv2.connectedComponentsWithStats(mask, 8)
    icons: list[VisualBlock] = []
    max_area = w * h * 0.08
    min_area = max(35, int(w * h * 0.00012))
    for label in range(1, num_labels):
        x, y, bw, bh, area = stats[label]
        rect = Rect(int(x), int(y), int(bw), int(bh)).padded(4, w, h)
        aspect = rect.w / max(1, rect.h)
        if area < min_area or rect.area > max_area or rect.w < 6 or rect.h < 6:
            continue
        if aspect > 7.0 or aspect < 0.12:
            continue
        if _is_mostly_white_region(arr, rect):
            continue
        image = _alpha_crop_for_content(slide_image, rect, exclude_rects)
        if image.mode == "RGBA" and np.count_nonzero(np.array(image)[:, :, 3]) < 18:
            continue
        icons.append(VisualBlock(rect, image, "icon"))

    return _dedupe_visuals(icons, w, h)


def _is_photo_like(crop: np.ndarray) -> bool:
    if crop.size == 0 or crop.shape[0] < 24 or crop.shape[1] < 24:
        return False

    gray = cv2.cvtColor(crop, cv2.COLOR_RGB2GRAY)
    variance = float(np.var(gray))
    edge_density = float(np.mean(cv2.Canny(gray, 60, 160) > 0))
    sample_w = min(64, crop.shape[1])
    sample_h = min(64, crop.shape[0])
    small = cv2.resize(crop, (sample_w, sample_h))
    unique_colors = len(np.unique(small.reshape(-1, 3), axis=0))
    return variance > 450 and edge_density > 0.035 and unique_colors > min(420, sample_w * sample_h * 0.18)


def _trim_content_rect(mask: np.ndarray, offset_x: int, offset_y: int, max_w: int, max_h: int) -> Rect | None:
    coords = cv2.findNonZero(mask.astype(np.uint8))
    if coords is None:
        return None
    x, y, w, h = cv2.boundingRect(coords)
    return Rect(int(x + offset_x), int(y + offset_y), int(w), int(h)).padded(2, max_w, max_h)


def _detect_layout_photo_regions(slide_image: Image.Image, text_rects: list[Rect]) -> list[VisualBlock]:
    arr = np.array(slide_image.convert("RGB"))
    h, w = arr.shape[:2]
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(arr, cv2.COLOR_RGB2HSV)
    saturation = hsv[:, :, 1]
    text_mask = _build_text_mask(w, h, text_rects, padding=8)
    content_mask = ((gray < 246) | (saturation > 24)).astype(np.uint8) * 255
    content_mask[text_mask > 0] = 0

    candidates: list[Rect] = []
    zones = [
        ("right", Rect(int(w * 0.48), 0, w - int(w * 0.48), h)),
        ("top", Rect(0, 0, w, int(h * 0.66))),
        ("upper_right", Rect(int(w * 0.52), 0, w - int(w * 0.52), int(h * 0.72))),
    ]
    for zone_name, zone in zones:
        zone_mask = content_mask[zone.y : zone.y2, zone.x : zone.x2]
        rect = _trim_content_rect(zone_mask, zone.x, zone.y, w, h)
        if rect is None:
            continue
        area_ratio = rect.area / max(1, w * h)
        if area_ratio < 0.07 or rect.w < w * 0.18 or rect.h < h * 0.18:
            continue
        if zone_name == "top":
            left_text_in_zone = sum(
                1 for text_rect in text_rects if text_rect.x < w * 0.48 and _intersection_area(text_rect, rect) > 0
            )
            if rect.w > w * 0.86 and left_text_in_zone >= 2:
                continue
        crop = arr[rect.y : rect.y2, rect.x : rect.x2]
        if _is_photo_like(crop):
            candidates.append(rect)

    photos: list[VisualBlock] = []
    for rect in _merge_rects(candidates, overlap_threshold=0.25):
        crop = slide_image.crop((rect.x, rect.y, rect.x2, rect.y2))
        photos.append(VisualBlock(rect=rect, image=crop, kind="photo"))
    return _dedupe_visuals(photos, w, h)


def _detect_magenta_icons(slide_image: Image.Image, exclude_rects: list[Rect]) -> list[VisualBlock]:
    arr = np.array(slide_image.convert("RGB"))
    h, w = arr.shape[:2]
    hsv = cv2.cvtColor(arr, cv2.COLOR_RGB2HSV)

    magenta_red = cv2.bitwise_or(
        cv2.inRange(hsv, np.array([145, 45, 35]), np.array([180, 255, 255])),
        cv2.inRange(hsv, np.array([0, 45, 35]), np.array([16, 255, 255])),
    )
    magenta_red = cv2.morphologyEx(magenta_red, cv2.MORPH_CLOSE, np.ones((3, 3), np.uint8), iterations=1)

    contours, _ = cv2.findContours(magenta_red, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    candidates: list[Rect] = []
    for contour in contours:
        x, y, bw, bh = cv2.boundingRect(contour)
        rect = Rect(int(x), int(y), int(bw), int(bh)).padded(6, w, h)
        area_ratio = rect.area / max(1, w * h)
        aspect = rect.w / max(1, rect.h)
        if rect.area < 50 or area_ratio > 0.08 or aspect > 6.0 or aspect < 0.15:
            continue
        if any(_overlap_ratio(rect, text_rect) > 0.35 for text_rect in exclude_rects):
            continue
        candidates.append(rect)

    icons: list[VisualBlock] = []
    for rect in _merge_rects(candidates, overlap_threshold=0.08):
        crop = arr[rect.y : rect.y2, rect.x : rect.x2]
        if crop.size == 0:
            continue

        crop_hsv = cv2.cvtColor(crop, cv2.COLOR_RGB2HSV)
        crop_gray = cv2.cvtColor(crop, cv2.COLOR_RGB2GRAY)
        alpha = cv2.bitwise_or(
            cv2.inRange(crop_hsv, np.array([145, 35, 25]), np.array([180, 255, 255])),
            cv2.inRange(crop_hsv, np.array([0, 35, 25]), np.array([16, 255, 255])),
        )
        alpha = cv2.bitwise_or(alpha, cv2.inRange(crop_gray, 0, 78))
        alpha = cv2.morphologyEx(alpha, cv2.MORPH_CLOSE, np.ones((2, 2), np.uint8), iterations=1)
        if np.count_nonzero(alpha) < 20:
            continue

        rgba = np.dstack([crop, alpha])
        icons.append(VisualBlock(rect=rect, image=Image.fromarray(rgba), kind="icon"))

    return _dedupe_visuals(icons, w, h)


def _classify_visual(arr: np.ndarray, rect: Rect) -> str:
    crop = arr[rect.y : rect.y2, rect.x : rect.x2]
    if crop.size == 0:
        return "visual"

    gray = cv2.cvtColor(crop, cv2.COLOR_RGB2GRAY)
    color_std = float(np.mean(np.std(crop.reshape(-1, 3), axis=0)))
    edge_density = float(np.mean(cv2.Canny(gray, 45, 130) > 0))
    if (
        rect.area > arr.shape[0] * arr.shape[1] * 0.04
        and rect.w > arr.shape[1] * 0.16
        and rect.h > arr.shape[0] * 0.16
        and (color_std > 34 or edge_density > 0.06)
    ):
        return "photo"
    if rect.area < arr.shape[0] * arr.shape[1] * 0.025 and (color_std > 18 or edge_density > 0.04):
        return "logo"
    return "visual"


def _region_complexity_score(crop: np.ndarray) -> tuple[float, float, list[str]]:
    if crop.size == 0:
        return 0.0, 1.0, ["empty crop"]

    gray = cv2.cvtColor(crop, cv2.COLOR_RGB2GRAY)
    edge_density = float(np.mean(cv2.Canny(gray, 50, 150) > 0))
    sample_w = min(80, crop.shape[1])
    sample_h = min(80, crop.shape[0])
    small = cv2.resize(crop, (sample_w, sample_h), interpolation=cv2.INTER_AREA)
    unique_ratio = len(np.unique(small.reshape(-1, 3), axis=0)) / max(1, sample_w * sample_h)
    color_std = float(np.mean(np.std(crop.reshape(-1, 3), axis=0))) / 128.0
    hist = cv2.calcHist([gray], [0], None, [32], [0, 256]).ravel()
    probabilities = hist / max(1.0, float(hist.sum()))
    entropy = float(-np.sum(probabilities[probabilities > 0] * np.log2(probabilities[probabilities > 0])) / 5.0)

    vector_score = 1.0 - min(1.0, unique_ratio * 1.7 + color_std * 0.45 + entropy * 0.22)
    if 0.01 <= edge_density <= 0.14:
        vector_score += 0.18
    elif edge_density > 0.24:
        vector_score -= 0.18
    vector_score = max(0.0, min(1.0, vector_score))

    raster_score = max(0.0, min(1.0, unique_ratio * 1.25 + color_std * 0.45 + entropy * 0.28 + max(0.0, edge_density - 0.12)))
    reasons = [
        f"unique_ratio={unique_ratio:.3f}",
        f"color_std={color_std:.3f}",
        f"entropy={entropy:.3f}",
        f"edge_density={edge_density:.3f}",
    ]
    return vector_score, raster_score, reasons


def _dominant_region_color(crop: np.ndarray) -> RGBColor:
    if crop.size == 0:
        return RGBColor(255, 255, 255)
    pixels = crop.reshape(-1, 3)
    median = np.median(pixels, axis=0)
    return RGBColor(int(median[0]), int(median[1]), int(median[2]))


def _detect_vector_elements(
    slide_image: Image.Image,
    text_rects: list[Rect],
    config: ReconstructionConfig,
) -> list[VectorElement]:
    arr = np.array(slide_image.convert("RGB"))
    h, w = arr.shape[:2]
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(arr, cv2.COLOR_RGB2HSV)
    saturation = hsv[:, :, 1]
    text_mask = _build_text_mask(w, h, text_rects, padding=7)

    content = ((gray < 246) | (saturation > 28)).astype(np.uint8) * 255
    content[text_mask > 0] = 0
    content = cv2.morphologyEx(content, cv2.MORPH_CLOSE, np.ones((3, 3), np.uint8), iterations=1)

    vectors: list[VectorElement] = []
    min_area = max(80, int(w * h * 0.00035))
    contours, _ = cv2.findContours(content, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    for contour in contours:
        x, y, bw, bh = cv2.boundingRect(contour)
        rect = Rect(int(x), int(y), int(bw), int(bh)).padded(1, w, h)
        if rect.area < min_area or rect.w < 4 or rect.h < 4:
            continue
        crop = arr[rect.y : rect.y2, rect.x : rect.x2]
        vector_score, raster_score, reasons = _region_complexity_score(crop)
        if raster_score >= config.raster_score_threshold and vector_score < config.vector_score_threshold:
            continue

        area = cv2.contourArea(contour)
        extent = area / max(1, bw * bh)
        peri = cv2.arcLength(contour, True)
        approx = cv2.approxPolyDP(contour, 0.025 * peri, True)
        fill = _dominant_region_color(crop)

        if rect.w > w * 0.92 and rect.h > h * 0.92:
            continue
        if len(approx) == 4 and extent > 0.72 and vector_score >= config.vector_score_threshold - 0.08:
            vectors.append(
                VectorElement("rect", rect, min(1.0, vector_score + 0.12), fill=fill, stroke=None, reasons=reasons + ["rectangular contour"])
            )
            continue

        aspect = rect.w / max(1, rect.h)
        if 0.65 <= aspect <= 1.55 and len(approx) >= 6 and extent > 0.58 and vector_score >= config.vector_score_threshold:
            vectors.append(VectorElement("ellipse", rect, vector_score, fill=fill, stroke=None, reasons=reasons + ["round contour"]))
            continue

        if len(approx) <= 8 and vector_score >= config.vector_score_threshold + 0.08:
            pts = [(int(point[0][0]), int(point[0][1])) for point in approx]
            vectors.append(VectorElement("polygon", rect, vector_score, fill=fill, stroke=None, points=pts, reasons=reasons + ["simple polygon"]))

    lines = cv2.HoughLinesP(content, 1, np.pi / 180, threshold=max(30, min(w, h) // 12), minLineLength=max(24, min(w, h) // 12), maxLineGap=8)
    if lines is not None:
        for line in lines[:80]:
            x1, y1, x2, y2 = [int(value) for value in line[0]]
            rect = Rect(min(x1, x2), min(y1, y2), abs(x2 - x1) + 1, abs(y2 - y1) + 1).padded(2, w, h)
            if any(_overlap_ratio(rect, text_rect) > 0.25 for text_rect in text_rects):
                continue
            length = ((x2 - x1) ** 2 + (y2 - y1) ** 2) ** 0.5
            if length < 24:
                continue
            crop = arr[rect.y : rect.y2, rect.x : rect.x2]
            vector_score, _raster_score, reasons = _region_complexity_score(crop)
            stroke = _dominant_region_color(crop)
            vectors.append(
                VectorElement("line", rect, max(vector_score, 0.72), stroke=stroke, stroke_width=2, points=[(x1, y1), (x2, y2)], reasons=reasons + ["hough line"])
            )

    return _dedupe_vectors(vectors)


def _dedupe_vectors(vectors: list[VectorElement]) -> list[VectorElement]:
    kept: list[VectorElement] = []
    for candidate in sorted(vectors, key=lambda item: (item.kind != "line", -item.score, -item.rect.area)):
        threshold = 0.72 if candidate.kind == "line" else 0.82
        if any(_overlap_ratio(candidate.rect, existing.rect) > threshold for existing in kept):
            continue
        kept.append(candidate)
    kept.sort(key=lambda item: (item.rect.y, item.rect.x))
    return kept


def _vector_to_svg_element(vector: VectorElement) -> str:
    fill = _rgb_to_hex(vector.fill) if vector.fill is not None else "none"
    stroke = _rgb_to_hex(vector.stroke) if vector.stroke is not None else "none"
    stroke_width = max(1, int(vector.stroke_width))
    if vector.kind == "rect":
        return f'<rect x="{vector.rect.x}" y="{vector.rect.y}" width="{vector.rect.w}" height="{vector.rect.h}" fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width}" />'
    if vector.kind == "ellipse":
        cx = vector.rect.x + vector.rect.w / 2
        cy = vector.rect.y + vector.rect.h / 2
        return f'<ellipse cx="{cx:.1f}" cy="{cy:.1f}" rx="{vector.rect.w / 2:.1f}" ry="{vector.rect.h / 2:.1f}" fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width}" />'
    if vector.kind == "line" and vector.points and len(vector.points) >= 2:
        (x1, y1), (x2, y2) = vector.points[:2]
        return f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" fill="none" stroke="{stroke}" stroke-width="{stroke_width}" stroke-linecap="round" />'
    if vector.points:
        points = " ".join(f"{x},{y}" for x, y in vector.points)
        return f'<polygon points="{points}" fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width}" />'
    return ""


def build_slide_svg(analyzed: AnalyzedSlide) -> str:
    elements = [_vector_to_svg_element(vector) for vector in (analyzed.vectors or [])]
    elements = [element for element in elements if element]
    body = "\n  ".join(elements)
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{analyzed.image.width}" height="{analyzed.image.height}" '
        f'viewBox="0 0 {analyzed.image.width} {analyzed.image.height}">\n  {body}\n</svg>\n'
    )


def _export_detection_overlay(source: Image.Image, rects: list[Rect], debug_dir: Path) -> None:
    overlay = source.copy().convert("RGB")
    draw = ImageDraw.Draw(overlay)
    for index, rect in enumerate(rects, start=1):
        draw.rectangle((rect.x, rect.y, rect.x2, rect.y2), outline=(255, 0, 0), width=4)
        draw.text((rect.x + 8, rect.y + 8), str(index), fill=(255, 0, 0))
    overlay.save(debug_dir / "erkannte_slide_grenzen.png")


def _draw_rects(image: Image.Image, rects: list[tuple[Rect, tuple[int, int, int]]]) -> Image.Image:
    overlay = image.copy().convert("RGB")
    draw = ImageDraw.Draw(overlay)
    for rect, color in rects:
        draw.rectangle((rect.x, rect.y, rect.x2, rect.y2), outline=color, width=2)
    return overlay


def _export_slide_debug(debug_dir: Path, index: int, analyzed: AnalyzedSlide) -> None:
    analyzed.image.save(debug_dir / f"slide_{index:03d}_crop.png")
    analyzed.image.save(debug_dir / f"slide_{index:03d}_without_border.png")
    vector_rects = [(vector.rect, (0, 180, 80)) for vector in analyzed.vectors or []]
    raster_rects = [(visual.rect, (235, 128, 0)) for visual in analyzed.visuals if visual.kind not in {"icon", "logo"}]
    icon_rects = [(visual.rect, (220, 30, 170)) for visual in analyzed.visuals if visual.kind == "icon"]
    text_rects = [(block.rect, (40, 80, 230)) for block in analyzed.texts]
    shape_rects = [(shape.rect, (150, 60, 190)) for shape in analyzed.shapes]
    _draw_rects(analyzed.image, vector_rects + raster_rects + icon_rects + text_rects + shape_rects).save(debug_dir / f"slide_{index:03d}_elements_overlay.png")
    _draw_rects(analyzed.image, vector_rects).save(debug_dir / f"slide_{index:03d}_vector_candidates.png")
    _draw_rects(analyzed.image, raster_rects).save(debug_dir / f"slide_{index:03d}_raster_fallbacks.png")
    asset_dir = debug_dir / "assets"
    asset_dir.mkdir(exist_ok=True)
    for visual_index, visual in enumerate(analyzed.visuals, start=1):
        if visual.kind in {"photo", "icon", "logo"}:
            visual.image.save(asset_dir / f"slide_{index:03d}_{visual.kind}_{visual_index:03d}.png")
    if analyzed.vectors is not None:
        (debug_dir / f"slide_{index:03d}.svg").write_text(build_slide_svg(analyzed), encoding="utf-8")
    (debug_dir / f"slide_{index:03d}_report.json").write_text(json.dumps(analyzed.debug or {}, indent=2), encoding="utf-8")


def _write_reconstruction_report(debug_dir: Path, slides: list[AnalyzedSlide]) -> None:
    report = {
        "slide_count": len(slides),
        "slides": [slide.debug or {} for slide in slides],
        "warnings": [],
        "errors": [],
    }
    (debug_dir / "reconstruction_report.json").write_text(json.dumps(report, indent=2), encoding="utf-8")


def _dedupe_visuals(visuals: list[VisualBlock], max_w: int, max_h: int) -> list[VisualBlock]:
    kept: list[VisualBlock] = []
    for candidate in sorted(visuals, key=lambda item: item.rect.area, reverse=True):
        if any(_overlap_ratio(candidate.rect, existing.rect) > 0.82 for existing in kept):
            continue
        kept.append(candidate)

    kept.sort(key=lambda item: (item.rect.y, item.rect.x))
    return kept


def _dedupe_shapes(shapes: list[ShapeBlock], max_w: int, max_h: int) -> list[ShapeBlock]:
    kept: list[ShapeBlock] = []
    for candidate in sorted(shapes, key=lambda item: item.rect.area, reverse=True):
        if candidate.rect.area > max_w * max_h * 0.86:
            continue
        if any(_overlap_ratio(candidate.rect, existing.rect) > 0.85 for existing in kept):
            continue
        kept.append(candidate)

    kept.sort(key=lambda item: (item.rect.y, item.rect.x))
    return kept


def _is_bullet_line(text: str) -> bool:
    stripped = text.strip()
    return stripped.startswith(("-", "*", "\u2022", "\u00b7", "\u2013", ">>", "\u00bb"))


def _clean_line_text(text: str) -> tuple[str, bool]:
    stripped = " ".join(text.split())
    is_bullet = _is_bullet_line(stripped)
    if is_bullet:
        stripped = stripped.lstrip("-*\u2022\u00b7\u2013\u00bb> ").strip()
    stripped = _strip_logo_noise(stripped)
    return stripped, is_bullet


def _polish_ocr_text(text: str) -> str:
    text = " ".join(text.split())
    replacements = {
        "Kl-gestutzteEntscheldungsintelligenz": "KI-gestützte Entscheidungsintelligenz",
        "Kl-gestutzte Entscheidungsintelligenz": "KI-gestützte Entscheidungsintelligenz",
        "Kl-gestutzte Entscheldungsintelligenz": "KI-gestützte Entscheidungsintelligenz",
        "inSRP-Processen": "in SAP-Prozessen",
        "in SRP-Processen": "in SAP-Prozessen",
        "Naturliche Sprache,": "Natürliche Sprache,",
        "Detenzusammenfassung,Vorharsagen": "Datenzusammenfassung, Vorhersagen",
        "Detenzusammenfassung, Vorharsagen": "Datenzusammenfassung, Vorhersagen",
        "Vurfugbar in SAP S/4HANA (Cloud),": "Verfügbar in SAP S/4HANA (Cloud),",
        "Vurfugbar in SAP S/4HANA (Cloud)": "Verfügbar in SAP S/4HANA (Cloud)",
        "insbesondereinSD-Processen": "insbesondere in SD-Prozessen",
        "SD-Processen": "SD-Prozessen",
        "Vertriebsunterstutzung": "Vertriebsunterstützung",
        "g mit SAP S/4HANA": "mit SAP S/4HANA",
        "ermoglicht": "ermöglicht",
        "fur": "für",
        "Skallerung": "Skalierung",
        "Serviceunterstutzung": "Serviceunterstützung",
        "SchnellereReaktion aufMarkt-und": "Schnellere Reaktion auf Markt- und",
        "BessereEntscheidungendurch": "Bessere Entscheidungen durch",
        "wenigermanaellenAnelyesoufwand": "weniger manuellen Analyseaufwand",
        "wenigermanualienAnalyse": "weniger manuellen Analyse",
        "WenigerSuchen,wenigerKlicken,": "Weniger Suchen, weniger Klicken,",
        "mehrEntschelden": "mehr Entscheiden",
        "Migrationsargumentation": "Migrationargumentation",
        "R/3 > S/4": "R/3 → S/4",
        "SAPR/3": "SAP R/3",
        "kompatihei": "kompatibel",
        "Jouleistverfugbarundwird": "Joule ist verfügbar und wird",
        "hauptsachlichuberFlori,BTPund": "hauptsächlich über Fiori, BTP und",
        "Standard-SD-Processe": "Standard-SD-Prozesse",
        "BassereEntscheidungen,weniger": "Bessere Entscheidungen, weniger",
        "manuellerAnalyseoufwand": "manueller Analyseaufwand",
        "KlassischeEmbeddedAnalytics": "Klassische Embedded Analytics",
        "Standard-KPls": "Standard-KPIs",
        "ErfahrenePower-User": "Erfahrene Power-User",
        "SkalierungvonExpertenwissen": "Skalierung von Expertenwissen",
    }
    for wrong, correct in replacements.items():
        text = text.replace(wrong, correct)
    return text


def _normalized_text(text: str) -> str:
    return "".join(ch.lower() for ch in text if ch.isalnum())


def _strip_logo_noise(text: str) -> str:
    stripped = text.strip()
    lowered = stripped.lower()
    for suffix in (".msg", " msg", "-msg", ".msq", " msq", ".6sw", " 6sw"):
        if lowered.endswith(suffix) and len(stripped) > len(suffix) + 6:
            return stripped[: -len(suffix)].rstrip(" .,-")
    return stripped


def _looks_like_logo_text(text: str) -> bool:
    normalized = _normalized_text(text)
    known = {"msg", "msq", "6sw", "sap", "ai"}
    return normalized in known or (2 <= len(normalized) <= 4 and any(ch.isalpha() for ch in normalized))


def _line_is_in_logo_zone(line: TextLine, width: int, height: int) -> bool:
    top_band = line.rect.y < height * 0.18
    corner_band = line.rect.x < width * 0.22 or line.rect.x2 > width * 0.72
    compact = line.rect.w < width * 0.22 and line.rect.h < height * 0.16
    return top_band and corner_band and compact


def _detect_logo_regions(slide_image: Image.Image, lines: list[TextLine]) -> list[VisualBlock]:
    if not lines:
        return []

    arr = np.array(slide_image.convert("RGB"))
    h, w = arr.shape[:2]
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(arr, cv2.COLOR_RGB2HSV)
    saturation = hsv[:, :, 1]
    content = ((gray < 245) | (saturation > 32)).astype(np.uint8) * 255
    content = cv2.morphologyEx(content, cv2.MORPH_CLOSE, np.ones((5, 5), np.uint8), iterations=1)
    num_labels, labels, stats, _centroids = cv2.connectedComponentsWithStats(content, 8)

    logos: list[VisualBlock] = []
    for line in lines:
        if not (_looks_like_logo_text(line.text) and _line_is_in_logo_zone(line, w, h)):
            continue

        seed = line.rect.padded(max(8, line.rect.h), w, h)
        seed_labels = labels[seed.y : seed.y2, seed.x : seed.x2]
        label_ids, counts = np.unique(seed_labels[seed_labels > 0], return_counts=True)
        rect = seed
        if label_ids.size:
            label = int(label_ids[int(np.argmax(counts))])
            x, y, bw, bh, area = stats[label]
            component = Rect(int(x), int(y), int(bw), int(bh)).padded(4, w, h)
            if 20 <= area <= w * h * 0.08 and component.w <= w * 0.34 and component.h <= h * 0.22:
                rect = component

        crop = slide_image.crop((rect.x, rect.y, rect.x2, rect.y2))
        logos.append(VisualBlock(rect=rect, image=crop, kind="logo"))

    return _dedupe_visuals(logos, w, h)


def _line_overlaps_protected_visual(line: TextLine, protected_visuals: list[VisualBlock]) -> bool:
    return any(_overlap_ratio(line.rect, visual.rect) > 0.72 for visual in protected_visuals)


def _has_bullet_marker(arr: np.ndarray, rect: Rect) -> bool:
    h, w = arr.shape[:2]
    x1 = max(0, rect.x - max(10, rect.h))
    x2 = max(0, rect.x - 2)
    y1 = max(0, rect.y - rect.h // 3)
    y2 = min(h, rect.y2 + rect.h // 3)
    if x2 <= x1 or y2 <= y1:
        return False

    crop = arr[y1:y2, x1:x2]
    gray = cv2.cvtColor(crop, cv2.COLOR_RGB2GRAY)
    hsv = cv2.cvtColor(crop, cv2.COLOR_RGB2HSV)
    sat = hsv[:, :, 1]
    mask = ((gray < 170) | ((gray < 230) & (sat > 45))).astype(np.uint8) * 255
    mask = cv2.morphologyEx(mask, cv2.MORPH_OPEN, np.ones((2, 2), np.uint8), iterations=1)
    num_labels, _labels, stats, _centroids = cv2.connectedComponentsWithStats(mask, 8)

    max_area = max(12, rect.h * rect.h)
    for label in range(1, num_labels):
        _x, _y, bw, bh, area = stats[label]
        if 4 <= area <= max_area and bw <= rect.h * 1.2 and bh <= rect.h * 1.2:
            return True
    return False


def _lines_belong_together(previous: TextLine, current: TextLine) -> bool:
    avg_h = max(1, (previous.rect.h + current.rect.h) / 2)
    vertical_gap = current.rect.y - previous.rect.y2
    if vertical_gap < -avg_h * 0.55 or vertical_gap > avg_h * 0.85:
        return False

    x_overlap = _intersection_area(
        Rect(previous.rect.x, 0, previous.rect.w, 1), Rect(current.rect.x, 0, current.rect.w, 1)
    )
    overlap_ratio = x_overlap / max(1, min(previous.rect.w, current.rect.w))
    left_delta = abs(previous.rect.x - current.rect.x)
    center_delta = abs((previous.rect.x + previous.rect.x2) / 2 - (current.rect.x + current.rect.x2) / 2)

    return overlap_ratio > 0.2 or left_delta <= avg_h * 3.2 or center_delta <= avg_h * 5.0


def _sort_text_blocks_reading_order(blocks: list[TextBlock], slide_width: int) -> list[TextBlock]:
    if not blocks:
        return []

    remaining = sorted(blocks, key=lambda block: (block.rect.y, block.rect.x))
    rows: list[list[TextBlock]] = []
    for block in remaining:
        center_y = block.rect.y + block.rect.h / 2
        for row in rows:
            row_center = sum(item.rect.y + item.rect.h / 2 for item in row) / len(row)
            tolerance = max(14, min(slide_width * 0.018, max(item.rect.h for item in row) * 0.75))
            if abs(center_y - row_center) <= tolerance:
                row.append(block)
                break
        else:
            rows.append([block])

    ordered: list[TextBlock] = []
    for row in sorted(rows, key=lambda items: min(item.rect.y for item in items)):
        ordered.extend(sorted(row, key=lambda item: item.rect.x))

    for index, block in enumerate(ordered):
        block.order = index
    return ordered


def _group_text_lines(lines: list[TextLine], slide_image: Image.Image) -> list[TextBlock]:
    if not lines:
        return []

    max_w, max_h = slide_image.width, slide_image.height
    arr = np.array(slide_image.convert("RGB"))
    sorted_lines = sorted(lines, key=lambda line: (line.rect.y, line.rect.x))
    prepared_lines: list[tuple[TextLine, str, bool]] = []
    for line in sorted_lines:
        cleaned, is_bullet = _clean_line_text(line.text)
        is_bullet = is_bullet or _has_bullet_marker(arr, line.rect)
        if cleaned:
            prepared_lines.append((line, cleaned, is_bullet))

    groups: list[list[tuple[TextLine, str, bool]]] = []
    for line, cleaned, is_bullet in prepared_lines:
        starts_new_bullet = bool(groups and is_bullet)
        if groups and not starts_new_bullet and _lines_belong_together(groups[-1][-1][0], line):
            groups[-1].append((line, cleaned, is_bullet))
        else:
            groups.append([(line, cleaned, is_bullet)])

    blocks: list[TextBlock] = []
    for group in groups:
        cleaned_lines: list[tuple[TextLine, bool]] = []
        for line, cleaned, is_bullet in group:
            cleaned_lines.append((TextLine(line.rect, _polish_ocr_text(cleaned), line.confidence, line.color), is_bullet))
        if not cleaned_lines:
            continue

        block_lines: list[TextLine] = []
        bullet_flags: list[bool] = []
        for line, is_bullet in cleaned_lines:
            block_lines.append(line)
            bullet_flags.append(is_bullet)

        rect = _union_rect([line.rect for line in block_lines], max_w, max_h, padding=2)
        blocks.append(TextBlock(rect, block_lines, bullet_flags))

    return _sort_text_blocks_reading_order(blocks, max_w)


def analyze_image(
    image_path: Path,
    confidence: float = MIN_TEXT_CONFIDENCE,
    include_visuals: bool = True,
    progress: Callable[[str], None] | None = None,
    config: ReconstructionConfig | None = None,
    debug_dir: Path | None = None,
) -> list[AnalyzedSlide]:
    progress = progress or (lambda _: None)
    config = config or DEFAULT_CONFIG
    source = Image.open(image_path).convert("RGB")
    if debug_dir is not None:
        debug_dir.mkdir(parents=True, exist_ok=True)
        source.save(debug_dir / "original_input.png")
    slide_rects = detect_slide_rects(source, config)
    progress(f"{len(slide_rects)} Folie(n) erkannt.")

    ocr = RapidOCR()
    slides: list[AnalyzedSlide] = []
    for index, rect in enumerate(slide_rects, start=1):
        progress(f"Analysiere Folie {index}/{len(slide_rects)} ...")
        crop = _inner_slide_crop(rect, source.width, source.height, config)
        slide_img = source.crop((crop.inner_rect.x, crop.inner_rect.y, crop.inner_rect.x2, crop.inner_rect.y2))
        arr = np.array(slide_img)
        ocr_image = slide_img
        if OCR_IMAGE_SCALE > 1:
            ocr_image = slide_img.resize(
                (slide_img.width * OCR_IMAGE_SCALE, slide_img.height * OCR_IMAGE_SCALE),
                Image.Resampling.LANCZOS,
            )
        result, _ = ocr(np.array(ocr_image))
        lines: list[TextLine] = []

        for item in result or []:
            points, text, score = item
            score_float = float(score)
            cleaned = " ".join(str(text).split())
            if score_float < confidence or not cleaned:
                continue
            scaled_points = [[point[0] / OCR_IMAGE_SCALE, point[1] / OCR_IMAGE_SCALE] for point in points]
            text_rect = _box_to_rect(scaled_points, slide_img.width, slide_img.height)
            if text_rect.w < 4 or text_rect.h < 4:
                continue
            lines.append(
                TextLine(
                    rect=text_rect,
                    text=cleaned,
                    confidence=score_float,
                    color=_guess_text_color(arr, text_rect),
                )
            )

        protected_visuals = _detect_logo_regions(slide_img, lines)
        editable_lines = [line for line in lines if not _line_overlaps_protected_visual(line, protected_visuals)]
        texts = _group_text_lines(editable_lines, slide_img)
        text_rects = [line.rect for line in editable_lines]
        visual_candidates = _detect_visuals(slide_img, text_rects) if include_visuals else []
        large_images = _detect_large_image_regions(slide_img, text_rects, visual_candidates) if include_visuals else []
        protected_rects = [visual.rect for visual in large_images + protected_visuals]
        icons = _detect_magenta_icons(slide_img, text_rects + protected_rects)
        icons.extend(_detect_icon_regions(slide_img, text_rects + protected_rects))
        remaining_visuals = [
            visual
            for visual in visual_candidates
            if not any(_overlap_ratio(visual.rect, large.rect) > 0.45 for large in large_images)
            and visual.kind in {"logo", "icon"}
        ]
        visuals = _dedupe_visuals(large_images + icons + remaining_visuals + protected_visuals, slide_img.width, slide_img.height)
        shapes = _detect_text_containers(slide_img, texts, [line.rect for line in editable_lines])
        vectors = _detect_vector_elements(slide_img, [line.rect for line in editable_lines], config)
        visuals = [
            visual
            for visual in visuals
            if visual.kind in {"photo", "logo", "icon"} or not any(_overlap_ratio(visual.rect, vector.rect) > 0.62 for vector in vectors)
        ]
        background = _rect_fill_color(arr, Rect(0, 0, slide_img.width, slide_img.height))
        slide_debug = {
            "index": index,
            "outer_rect": crop.outer_rect.__dict__,
            "inner_rect": crop.inner_rect.__dict__,
            "border_trim": {
                "left": crop.inner_rect.x - crop.outer_rect.x,
                "top": crop.inner_rect.y - crop.outer_rect.y,
                "right": crop.outer_rect.x2 - crop.inner_rect.x2,
                "bottom": crop.outer_rect.y2 - crop.inner_rect.y2,
            },
            "texts": [
                {"text": block.text, "confidence": block.confidence, "order": block.order, "rect": block.rect.__dict__}
                for block in texts
            ],
            "visuals": [{"kind": visual.kind, "rect": visual.rect.__dict__} for visual in visuals],
            "shapes": [{"kind": "rect", "rect": shape.rect.__dict__, "fill": _rgb_to_tuple(shape.fill)} for shape in shapes],
            "vectors": [
                {
                    "kind": vector.kind,
                    "rect": vector.rect.__dict__,
                    "score": round(vector.score, 3),
                    "fill": _rgb_to_tuple(vector.fill),
                    "stroke": _rgb_to_tuple(vector.stroke),
                    "reasons": vector.reasons or [],
                }
                for vector in vectors
            ],
        }
        analyzed = AnalyzedSlide(
            rect=crop.inner_rect,
            image=slide_img,
            texts=texts,
            visuals=visuals,
            shapes=shapes,
            vectors=vectors,
            background=background,
            source_rect=crop.outer_rect,
            debug=slide_debug,
        )
        slides.append(analyzed)
        if debug_dir is not None:
            _export_slide_debug(debug_dir, index, analyzed)
        logo_count = sum(1 for visual in visuals if visual.kind == "logo")
        icon_count = sum(1 for visual in visuals if visual.kind == "icon")
        progress(
            f"Folie {index}: {len(texts)} Textfeld(er), {len(visuals)} Bildobjekt(e), "
            f"{len(shapes)} Flaeche(n), {len(vectors)} Vektor(en), {logo_count} Logo(s), {icon_count} Icon(s)."
        )

    if debug_dir is not None:
        _export_detection_overlay(source, slide_rects, debug_dir)
        _write_reconstruction_report(debug_dir, slides)
    return slides


def _add_textbox(slide, block: TextBlock, scale_x: float, scale_y: float) -> None:
    left = int(block.rect.x * scale_x)
    top = int(block.rect.y * scale_y)
    width = int(max(1, block.rect.w * scale_x * 1.12))
    height = int(max(1, block.rect.h * scale_y * 1.45))
    median_line_h = sorted(line.rect.h for line in block.lines)[len(block.lines) // 2]
    base_size = max(7, min(30, int(median_line_h * scale_y / EMU_PER_INCH * 72 * 1.05)))
    line_texts = [line.text for line in block.lines]
    font_size = Pt(_fit_font_size(line_texts, base_size, width, height))

    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.word_wrap = True
    box.text_frame.clear()

    for index, line in enumerate(block.lines):
        p = box.text_frame.paragraphs[0] if index == 0 else box.text_frame.add_paragraph()
        is_bullet = index < len(block.bullet_lines) and block.bullet_lines[index]
        p.text = f"\u2022 {line.text}" if is_bullet else line.text
        p.font.name = "Aptos"
        p.font.size = font_size
        p.font.color.rgb = line.color


def _resolve_text_block_overlaps(blocks: list[TextBlock], slide_w: int, slide_h: int) -> list[TextBlock]:
    ordered = _sort_text_blocks_reading_order(list(blocks), slide_w)
    resolved: list[TextBlock] = []
    for block in ordered:
        rect = block.rect
        for existing in resolved:
            if _intersection_area(rect, existing.rect) == 0:
                continue
            if rect.y >= existing.rect.y or rect.x < existing.rect.x2:
                rect = rect.moved(rect.x, existing.rect.y2 + 3, slide_w, slide_h)
        resolved.append(TextBlock(rect, block.lines, block.bullet_lines, block.order))
    return resolved


def _shape_rect_from_px(rect: Rect, scale_x: float, scale_y: float) -> tuple[int, int, int, int]:
    return (
        int(rect.x * scale_x),
        int(rect.y * scale_y),
        int(max(1, rect.w * scale_x)),
        int(max(1, rect.h * scale_y)),
    )


def _template_kind(index: int, analyzed: AnalyzedSlide) -> str:
    all_text = " ".join(block.text for block in analyzed.texts)
    normalized = _normalized_text(all_text)
    if "fazit" in normalized:
        return "conclusion"
    if "migration" in normalized or "r3" in normalized:
        return "migration"
    if "joulevsklassischeembeddedanalytics" in normalized:
        return "comparison"
    if "mehrwertvonjouleimsd" in normalized:
        return "card_grid"
    if index == 1:
        return "cover"
    if any(visual.kind == "photo" and visual.rect.x > analyzed.image.width * 0.45 for visual in analyzed.visuals):
        return "photo_right"
    return "standard"


def _is_styled_text_noise(block: TextBlock, template: str, slide_width_px: int, slide_height_px: int) -> bool:
    normalized = _normalized_text(block.text)
    if not normalized:
        return True
    if normalized in {"ai", "ct", "oo", "00"}:
        return True
    if block.rect.x > slide_width_px * 0.62 and len(normalized) <= 4:
        return True
    if template in {"photo_right", "standard"} and block.rect.x > slide_width_px * 0.62 and block.rect.y > slide_height_px * 0.18:
        return len(normalized) < 12
    return False


def _text_role(block: TextBlock, slide_width_px: int, slide_height_px: int, template: str = "standard") -> str:
    text = block.text.strip()
    normalized = _normalized_text(text)
    y_ratio = block.rect.y / max(1, slide_height_px)
    area_ratio = block.rect.area / max(1, slide_width_px * slide_height_px)
    has_bullet = any(block.bullet_lines)

    if normalized in {"msg", "msq", "6sw"}:
        return "logo"
    if template == "cover" and y_ratio > 0.5:
        return "cover_title"
    if has_bullet:
        return "bullet"
    if template == "conclusion" and "fazit" in normalized:
        return "heading"
    if y_ratio < 0.18 and block.rect.w > slide_width_px * 0.32:
        return "title"
    if template in {"card_grid", "comparison"} and y_ratio > 0.18 and (text[:2].isdigit() or len(text) < 45):
        return "card_heading"
    if y_ratio > 0.55 and area_ratio > 0.04 and len(block.lines) >= 2:
        return "cover_title"
    if len(block.lines) == 1 and (block.rect.h > slide_height_px * 0.045 or len(text) < 45):
        return "heading"
    if text[:2].isdigit() or text[:2] in {"1.", "2.", "3.", "4.", "5.", "6."}:
        return "card_heading"
    return "body"


def _fit_font_size(lines: list[str], base_size: int, width: int, height: int) -> int:
    if not lines:
        return base_size
    width_in = max(0.1, width / EMU_PER_INCH)
    height_in = max(0.1, height / EMU_PER_INCH)
    max_chars = max(len(line) for line in lines)
    line_count = max(1, len(lines))
    size = base_size
    if max_chars:
        estimated_chars_per_line = width_in * 72 / max(size * 0.52, 1)
        if estimated_chars_per_line < max_chars:
            size = min(size, int(max(6, size * estimated_chars_per_line / max_chars)))
    estimated_line_height = size * 1.18 / 72
    if estimated_line_height * line_count > height_in:
        size = min(size, int(max(6, height_in * 72 / (line_count * 1.18))))
    return max(6, min(base_size, size))


def _set_run_font(run, style: TextStyle) -> None:
    run.font.name = style.font_name
    run.font.size = Pt(style.size)
    run.font.bold = style.bold
    run.font.color.rgb = style.color


def _configure_text_frame(shape) -> None:
    shape.text_frame.margin_left = 0
    shape.text_frame.margin_right = 0
    shape.text_frame.margin_top = 0
    shape.text_frame.margin_bottom = 0
    shape.text_frame.word_wrap = True
    shape.text_frame.clear()


def _add_plain_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    lines: list[str],
    style: TextStyle | int,
    color: RGBColor | None = None,
    bold: bool = False,
    font_name: str = "Aptos",
) -> None:
    if isinstance(style, int):
        style = TextStyle(style, color or REFERENCE_TEXT, bold=bold, font_name=font_name)
    fitted_style = TextStyle(
        _fit_font_size(lines, style.size, width, height),
        style.color,
        bold=style.bold,
        font_name=style.font_name,
    )
    box = slide.shapes.add_textbox(left, top, width, height)
    _configure_text_frame(box)
    for index, line in enumerate(lines):
        line = line.replace("\u00e2\u20ac\u00a2", "\u2022")
        paragraph = box.text_frame.paragraphs[0] if index == 0 else box.text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        _set_run_font(run, fitted_style)


def _add_styled_textbox(
    slide, block: TextBlock, scale_x: float, scale_y: float, slide_w: int, slide_h: int, template: str = "standard"
) -> None:
    left, top, width, height = _shape_rect_from_px(block.rect, scale_x, scale_y)
    role = _text_role(block, slide_w, slide_h, template)
    lines = [_polish_ocr_text(line.text) for line in block.lines]

    if role == "logo":
        _add_plain_textbox(
            slide,
            left,
            top,
            max(width, int(slide_w * scale_x * 0.09)),
            max(height, int(slide_h * scale_y * 0.045)),
            lines,
            21,
            REFERENCE_LOGO,
            bold=True,
            font_name="Aptos Display",
        )
        return

    if role == "cover_title":
        title_lines = lines[:2]
        subtitle_lines = lines[2:]
        title_height = max(1, int(height * 0.62))
        _add_plain_textbox(
            slide,
            left,
            top,
            int(width * 0.72),
            title_height,
            title_lines,
            18,
            REFERENCE_ACCENT,
            bold=True,
            font_name="Aptos Display",
        )
        if subtitle_lines:
            _add_plain_textbox(
                slide,
                left,
                top + title_height,
                int(width * 0.68),
                max(1, height - title_height),
                subtitle_lines,
                10,
                REFERENCE_TEXT,
            )
        return

    if role == "bullet":
        bullet_size = 14 if block.rect.h > slide_h * 0.045 else 11
        text_size = 12 if block.rect.h > slide_h * 0.05 else 9
        marker_left = max(0, left - int(slide_w * scale_x * 0.02))
        marker_width = max(1, int(slide_w * scale_x * 0.018))
        _add_plain_textbox(
            slide,
            marker_left,
            top,
            marker_width,
            max(height, int(slide_h * scale_y * 0.035)),
            ["•"],
            bullet_size,
            REFERENCE_ACCENT,
            bold=True,
        )
        _add_plain_textbox(
            slide,
            left,
            top,
            int(width * 1.14),
            max(height, int(slide_h * scale_y * 0.07)),
            lines,
            text_size,
            REFERENCE_TEXT,
        )
        return

    if role == "title":
        _add_plain_textbox(
            slide,
            left,
            top,
            int(width * 1.18),
            max(height, int(slide_h * scale_y * 0.07)),
            lines,
            21 if block.rect.w < slide_w * 0.65 else 20,
            REFERENCE_TEXT,
            bold=True,
            font_name="Aptos Display",
        )
        return

    if role in {"heading", "card_heading"}:
        _add_plain_textbox(
            slide,
            left,
            top,
            int(width * 1.12),
            max(height, int(slide_h * scale_y * 0.055)),
            lines,
            8 if role == "card_heading" else 21,
            REFERENCE_TEXT,
            bold=True,
            font_name="Aptos Display" if role == "heading" else "Aptos",
        )
        return

    _add_plain_textbox(
        slide,
        left,
        top,
        int(width * 1.08),
        max(height, int(slide_h * scale_y * 0.06)),
        lines,
        9 if block.rect.w < slide_w * 0.34 else 12,
        REFERENCE_TEXT,
    )


def _add_image_to_slide(slide, image: Image.Image, image_path: Path, left: int, top: int, width: int, height: int) -> None:
    image.save(image_path)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def _set_slide_background(slide, color: RGBColor | None) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or RGBColor(255, 255, 255)


def _add_vector_element(slide, vector: VectorElement, scale_x: float, scale_y: float) -> None:
    if vector.kind == "line" and vector.points and len(vector.points) >= 2:
        (x1, y1), (x2, y2) = vector.points[:2]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            int(x1 * scale_x),
            int(y1 * scale_y),
            int(x2 * scale_x),
            int(y2 * scale_y),
        )
        connector.line.width = max(1, int(vector.stroke_width * min(scale_x, scale_y)))
        if vector.stroke is not None:
            connector.line.color.rgb = vector.stroke
        return

    left, top, width, height = _shape_rect_from_px(vector.rect, scale_x, scale_y)
    if vector.kind == "ellipse":
        shape_type = MSO_SHAPE.OVAL
    elif vector.kind == "polygon" and vector.points and len(vector.points) == 3:
        shape_type = MSO_SHAPE.ISOSCELES_TRIANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if vector.fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = vector.fill
    if vector.stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = vector.stroke


def _visual_is_raster_fallback(visual: VisualBlock, mode: str) -> bool:
    if mode == "vector":
        return visual.kind == "photo"
    return visual.kind not in {"icon", "logo"}


def validate_pptx(output_path: Path) -> tuple[bool, list[str]]:
    errors: list[str] = []
    if not output_path.exists():
        return False, ["PPTX-Datei existiert nicht."]
    try:
        with zipfile.ZipFile(output_path) as archive:
            names = set(archive.namelist())
            for required in ("[Content_Types].xml", "ppt/presentation.xml", "_rels/.rels"):
                if required not in names:
                    errors.append(f"Pflichtdatei fehlt: {required}")
            for name in names:
                if name.endswith(".xml") or name.endswith(".rels"):
                    try:
                        etree.fromstring(archive.read(name))
                    except Exception as exc:
                        errors.append(f"Ungueltiges XML in {name}: {exc}")
    except zipfile.BadZipFile as exc:
        errors.append(f"PPTX ist kein gueltiges ZIP: {exc}")
    return not errors, errors


def export_pptx(
    slides: list[AnalyzedSlide],
    output_path: Path,
    progress: Callable[[str], None] | None = None,
    mode: str = "visual_safe",
    config: ReconstructionConfig | None = None,
    debug_dir: Path | None = None,
) -> None:
    if not slides:
        raise ValueError("Keine Folien zum Exportieren gefunden.")
    if mode not in EXPORT_MODES:
        raise ValueError(f"Unbekannter Exportmodus: {mode}")

    config = config or DEFAULT_CONFIG
    if debug_dir is not None:
        debug_dir.mkdir(parents=True, exist_ok=True)
    progress = progress or (lambda _: None)
    progress(f"Exportmodus: {mode}.")
    prs = Presentation()
    prs.slide_width = int(config.slide_width_inches * EMU_PER_INCH)
    prs.slide_height = int(config.slide_height_inches * EMU_PER_INCH)
    blank_layout = prs.slide_layouts[6]
    effective_mode = "visual_safe" if mode == "pixel" else mode

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_root = Path(temp_dir)
        for index, analyzed in enumerate(slides, start=1):
            progress(f"Exportiere Folie {index}/{len(slides)} ...")
            slide = prs.slides.add_slide(blank_layout)
            template = _template_kind(index, analyzed) if effective_mode == "styled_reconstruct" else "standard"

            scale_x = prs.slide_width / analyzed.image.width
            scale_y = prs.slide_height / analyzed.image.height

            if effective_mode in {"visual_safe", "editable", "styled_reconstruct"}:
                background = analyzed.image
                if effective_mode in {"editable", "styled_reconstruct"}:
                    line_rects = [line.rect for block in analyzed.texts for line in block.lines]
                    background = _remove_text_from_image(analyzed.image, line_rects)
                bg_path = temp_root / f"slide_{index:03d}_background.png"
                _add_image_to_slide(slide, background, bg_path, 0, 0, int(prs.slide_width), int(prs.slide_height))
            else:
                _set_slide_background(slide, analyzed.background)

            if effective_mode == "reconstruct":
                export_visuals = analyzed.visuals
            elif mode in {"hybrid", "vector"}:
                export_visuals = [visual for visual in analyzed.visuals if _visual_is_raster_fallback(visual, mode)]
            else:
                export_visuals = [visual for visual in analyzed.visuals if visual.kind in {"icon", "logo"}]

            for visual_index, visual in enumerate(export_visuals, start=1):
                if not config.enable_raster_fallback and mode in {"hybrid", "vector"}:
                    continue
                image_path = temp_root / f"slide_{index:03d}_visual_{visual_index:03d}.png"
                _add_image_to_slide(
                    slide,
                    visual.image,
                    image_path,
                    int(visual.rect.x * scale_x),
                    int(visual.rect.y * scale_y),
                    int(visual.rect.w * scale_x),
                    int(visual.rect.h * scale_y),
                )

            if effective_mode == "reconstruct" or mode in {"hybrid", "vector"}:
                for shape in analyzed.shapes:
                    rect = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        int(shape.rect.x * scale_x),
                        int(shape.rect.y * scale_y),
                        int(shape.rect.w * scale_x),
                        int(shape.rect.h * scale_y),
                    )
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = shape.fill
                    rect.line.fill.background()

            if mode in {"hybrid", "vector"} and config.enable_native_shapes:
                for vector in analyzed.vectors or []:
                    _add_vector_element(slide, vector, scale_x, scale_y)
                if debug_dir is not None and config.enable_svg_export:
                    (debug_dir / f"slide_{index:03d}.svg").write_text(build_slide_svg(analyzed), encoding="utf-8")

            export_texts = _resolve_text_block_overlaps(analyzed.texts, analyzed.image.width, analyzed.image.height)
            for block in export_texts:
                if effective_mode == "styled_reconstruct":
                    if _is_styled_text_noise(block, template, analyzed.image.width, analyzed.image.height):
                        continue
                    _add_styled_textbox(
                        slide, block, scale_x, scale_y, analyzed.image.width, analyzed.image.height, template
                    )
                else:
                    _add_textbox(slide, block, scale_x, scale_y)

        prs.save(output_path)
    valid, errors = validate_pptx(output_path)
    if not valid:
        raise ValueError("Die erzeugte PPTX ist technisch ungueltig: " + "; ".join(errors[:5]))
    progress("PPTX-Validierung erfolgreich.")


def convert_image_to_pptx(
    image_path: Path,
    output_path: Path,
    confidence: float = MIN_TEXT_CONFIDENCE,
    include_visuals: bool = True,
    progress: Callable[[str], None] | None = None,
    mode: str = "visual_safe",
    config: ReconstructionConfig | None = None,
    debug_dir: Path | None = None,
) -> None:
    config = config or DEFAULT_CONFIG
    if debug_dir is None and config.enable_debug_output:
        debug_dir = output_path.parent / f"{output_path.stem}_debug"
    if debug_dir is not None:
        debug_dir.mkdir(parents=True, exist_ok=True)
    slides = analyze_image(image_path, confidence, include_visuals, progress, config, debug_dir)
    export_pptx(slides, output_path, progress, mode, config, debug_dir)


def configure_windows_taskbar_icon() -> None:
    if sys.platform != "win32":
        return
    try:
        ctypes.windll.shell32.SetCurrentProcessExplicitAppUserModelID(WINDOWS_APP_ID)
    except Exception:
        pass


def show_startup_splash(root: tk.Tk, minimum_ms: int = 3000) -> None:
    if os.environ.get("SNAP2SLIDES_SKIP_SPLASH") == "1" or not APP_SPLASH_PNG.exists():
        return

    root.withdraw()
    splash = tk.Toplevel(root)
    splash.overrideredirect(True)
    splash.configure(background="#020b1c")

    screen_w = splash.winfo_screenwidth()
    screen_h = splash.winfo_screenheight()
    max_w = min(980, int(screen_w * 0.82))
    max_h = min(560, int(screen_h * 0.72))

    image = Image.open(APP_SPLASH_PNG).convert("RGB")
    image.thumbnail((max_w, max_h), Image.Resampling.LANCZOS)
    splash._photo = ImageTk.PhotoImage(image)  # type: ignore[attr-defined]
    label = ttk.Label(splash, image=splash._photo, borderwidth=0)  # type: ignore[attr-defined]
    label.pack()

    splash.update_idletasks()
    x = (screen_w - splash.winfo_width()) // 2
    y = (screen_h - splash.winfo_height()) // 2
    splash.geometry(f"+{x}+{y}")
    splash.lift()

    def finish_splash() -> None:
        splash.destroy()
        root.deiconify()
        root.lift()

    root.after(minimum_ms, finish_splash)


def _move_text_block(block: TextBlock, rect: Rect) -> TextBlock:
    dx = rect.x - block.rect.x
    dy = rect.y - block.rect.y
    moved_lines = [
        TextLine(Rect(line.rect.x + dx, line.rect.y + dy, line.rect.w, line.rect.h), line.text, line.confidence, line.color)
        for line in block.lines
    ]
    return TextBlock(rect, moved_lines, list(block.bullet_lines), block.order)


def _set_text_block_text(block: TextBlock, text: str) -> TextBlock:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        lines = [""]
    line_h = max(6, block.rect.h // max(1, len(lines)))
    text_lines: list[TextLine] = []
    bullet_flags: list[bool] = []
    for index, line_text in enumerate(lines):
        cleaned, is_bullet = _clean_line_text(line_text)
        y = block.rect.y + min(block.rect.h - line_h, index * line_h)
        text_lines.append(TextLine(Rect(block.rect.x, y, block.rect.w, line_h), cleaned, block.confidence, block.color))
        bullet_flags.append(is_bullet)
    return TextBlock(block.rect, text_lines, bullet_flags, block.order)


class ReconstructionEditor(tk.Toplevel):
    def __init__(
        self,
        parent: tk.Tk,
        slides: list[AnalyzedSlide],
        on_export: Callable[[list[AnalyzedSlide]], None],
    ) -> None:
        super().__init__(parent)
        self.title("Folien-Rekonstruktion bearbeiten")
        self.geometry("1040x680")
        self.minsize(880, 560)
        self.transient(parent)
        self.slides = slides
        self.on_export = on_export
        self.slide_index = 0
        self.selected: tuple[str, int] | None = None
        self.drag_start: tuple[int, int, Rect] | None = None
        self.scale = 1.0
        self._canvas_photo: ImageTk.PhotoImage | None = None
        self._item_map: dict[int, tuple[str, int]] = {}

        self._build_editor_ui()
        self._render_slide()
        self.grab_set()

    def _build_editor_ui(self) -> None:
        self.columnconfigure(0, weight=1)
        self.rowconfigure(0, weight=1)
        main = ttk.Frame(self, padding=12)
        main.grid(row=0, column=0, sticky="nsew")
        main.columnconfigure(0, weight=1)
        main.rowconfigure(1, weight=1)

        toolbar = ttk.Frame(main)
        toolbar.grid(row=0, column=0, columnspan=2, sticky="ew", pady=(0, 8))
        ttk.Button(toolbar, text="Vorherige", command=lambda: self._change_slide(-1)).pack(side=tk.LEFT)
        ttk.Button(toolbar, text="Naechste", command=lambda: self._change_slide(1)).pack(side=tk.LEFT, padx=(6, 12))
        self.slide_label = ttk.Label(toolbar, text="")
        self.slide_label.pack(side=tk.LEFT)
        ttk.Button(toolbar, text="PPTX erzeugen", command=self._finish).pack(side=tk.RIGHT)

        self.canvas = tk.Canvas(main, background="#f4f4f4", highlightthickness=1, highlightbackground="#b7b7b7")
        self.canvas.grid(row=1, column=0, sticky="nsew")
        self.canvas.bind("<Button-1>", self._select_at)
        self.canvas.bind("<B1-Motion>", self._drag_selected)
        self.canvas.bind("<ButtonRelease-1>", lambda _event: self._stop_drag())

        side = ttk.Frame(main, width=260)
        side.grid(row=1, column=1, sticky="ns", padx=(12, 0))
        side.grid_propagate(False)
        ttk.Label(side, text="Text").pack(anchor="w")
        self.text_editor = tk.Text(side, height=9, wrap=tk.WORD)
        self.text_editor.pack(fill=tk.X, pady=(4, 8))
        ttk.Button(side, text="Text uebernehmen", command=self._apply_text).pack(fill=tk.X)
        ttk.Separator(side).pack(fill=tk.X, pady=10)
        ttk.Button(side, text="Groesser", command=lambda: self._scale_selected(1.08)).pack(fill=tk.X)
        ttk.Button(side, text="Kleiner", command=lambda: self._scale_selected(0.92)).pack(fill=tk.X, pady=(4, 0))
        ttk.Button(side, text="Bild ersetzen", command=self._replace_visual).pack(fill=tk.X, pady=(12, 0))
        self.selection_label = ttk.Label(side, text="Kein Objekt ausgewaehlt", wraplength=240)
        self.selection_label.pack(anchor="w", pady=(14, 0))

    @property
    def current_slide(self) -> AnalyzedSlide:
        return self.slides[self.slide_index]

    def _change_slide(self, delta: int) -> None:
        self.slide_index = max(0, min(len(self.slides) - 1, self.slide_index + delta))
        self.selected = None
        self._render_slide()

    def _render_slide(self) -> None:
        slide = self.current_slide
        self.canvas.delete("all")
        self._item_map.clear()
        canvas_w = max(360, self.canvas.winfo_width() or 720)
        canvas_h = max(240, self.canvas.winfo_height() or 420)
        self.scale = min((canvas_w - 20) / slide.image.width, (canvas_h - 20) / slide.image.height)
        preview = slide.image.copy()
        preview.thumbnail((int(slide.image.width * self.scale), int(slide.image.height * self.scale)))
        self._canvas_photo = ImageTk.PhotoImage(preview)
        self.canvas.create_image(10, 10, image=self._canvas_photo, anchor=tk.NW)

        for index, visual in enumerate(slide.visuals):
            color = "#d97904" if visual.kind == "photo" else "#c0267c"
            self._draw_object("visual", index, visual.rect, color)
        for index, block in enumerate(slide.texts):
            self._draw_object("text", index, block.rect, "#2756d8")
        self.slide_label.configure(text=f"Folie {self.slide_index + 1} von {len(self.slides)}")
        self._update_selection_panel()

    def _draw_object(self, kind: str, index: int, rect: Rect, color: str) -> None:
        x1 = 10 + rect.x * self.scale
        y1 = 10 + rect.y * self.scale
        x2 = 10 + rect.x2 * self.scale
        y2 = 10 + rect.y2 * self.scale
        width = 3 if self.selected == (kind, index) else 2
        item = self.canvas.create_rectangle(x1, y1, x2, y2, outline=color, width=width)
        self._item_map[item] = (kind, index)
        if kind == "text":
            text = self.current_slide.texts[index].text.replace("\n", " ")
            label = self.canvas.create_text(x1 + 4, y1 + 4, text=text[:44], anchor=tk.NW, fill=color, width=max(40, x2 - x1 - 8))
            self._item_map[label] = (kind, index)

    def _select_at(self, event: tk.Event) -> None:
        hits = self.canvas.find_overlapping(event.x - 2, event.y - 2, event.x + 2, event.y + 2)
        self.selected = None
        for item in reversed(hits):
            if item in self._item_map:
                self.selected = self._item_map[item]
                break
        if self.selected is not None:
            rect = self._selected_rect()
            if rect is not None:
                self.drag_start = (event.x, event.y, rect)
        self._render_slide()

    def _selected_rect(self) -> Rect | None:
        if self.selected is None:
            return None
        kind, index = self.selected
        if kind == "text" and index < len(self.current_slide.texts):
            return self.current_slide.texts[index].rect
        if kind == "visual" and index < len(self.current_slide.visuals):
            return self.current_slide.visuals[index].rect
        return None

    def _set_selected_rect(self, rect: Rect) -> None:
        if self.selected is None:
            return
        kind, index = self.selected
        slide = self.current_slide
        if kind == "text" and index < len(slide.texts):
            slide.texts[index] = _move_text_block(slide.texts[index], rect)
        elif kind == "visual" and index < len(slide.visuals):
            slide.visuals[index].rect = rect

    def _drag_selected(self, event: tk.Event) -> None:
        if self.drag_start is None:
            return
        start_x, start_y, rect = self.drag_start
        dx = int((event.x - start_x) / max(self.scale, 0.001))
        dy = int((event.y - start_y) / max(self.scale, 0.001))
        slide = self.current_slide
        self._set_selected_rect(rect.moved(rect.x + dx, rect.y + dy, slide.image.width, slide.image.height))
        self._render_slide()

    def _stop_drag(self) -> None:
        self.drag_start = None

    def _scale_selected(self, factor: float) -> None:
        rect = self._selected_rect()
        if rect is None:
            return
        slide = self.current_slide
        self._set_selected_rect(rect.resized(rect.w * factor, rect.h * factor, slide.image.width, slide.image.height))
        self._render_slide()

    def _update_selection_panel(self) -> None:
        self.text_editor.delete("1.0", tk.END)
        if self.selected is None:
            self.selection_label.configure(text="Kein Objekt ausgewaehlt")
            return
        kind, index = self.selected
        if kind == "text":
            block = self.current_slide.texts[index]
            self.text_editor.insert("1.0", block.text)
            self.selection_label.configure(text=f"Textblock {index + 1}, Reihenfolge {block.order + 1}")
        else:
            visual = self.current_slide.visuals[index]
            self.selection_label.configure(text=f"Bildobjekt {index + 1}: {visual.kind}")

    def _apply_text(self) -> None:
        if self.selected is None or self.selected[0] != "text":
            return
        index = self.selected[1]
        self.current_slide.texts[index] = _set_text_block_text(self.current_slide.texts[index], self.text_editor.get("1.0", tk.END))
        self._render_slide()

    def _replace_visual(self) -> None:
        if self.selected is None or self.selected[0] != "visual":
            return
        path = filedialog.askopenfilename(
            parent=self,
            title="Bild ersetzen",
            filetypes=[("Bilder", "*.png *.jpg *.jpeg *.bmp *.tif *.tiff"), ("Alle Dateien", "*.*")],
        )
        if not path:
            return
        index = self.selected[1]
        self.current_slide.visuals[index].image = Image.open(path).convert("RGBA")
        self._render_slide()

    def _finish(self) -> None:
        self.grab_release()
        self.destroy()
        self.on_export(self.slides)


class Image2PptApp(tk.Tk):
    def __init__(self) -> None:
        super().__init__()
        self.title("Image2PPTSlicer")
        self._set_app_icon()
        self.geometry("860x620")
        self.minsize(760, 520)

        self.input_path = tk.StringVar()
        self.output_path = tk.StringVar()
        self.confidence = tk.DoubleVar(value=MIN_TEXT_CONFIDENCE)
        self.confidence_label = tk.StringVar()
        self.include_visuals = tk.BooleanVar(value=True)
        self.export_mode = tk.StringVar(value="hybrid")
        self.status_text = tk.StringVar(value="Bereit")
        self._preview_image: ImageTk.PhotoImage | None = None
        self._preview_source: Image.Image | None = None
        self._is_running = False
        self._pending_output_path: Path | None = None
        self._pending_mode: str = "hybrid"
        self._pending_config: ReconstructionConfig = DEFAULT_CONFIG

        self._build_ui()
        self._bind_shortcuts()
        self.input_path.trace_add("write", lambda *_args: self._mark_output_stale())
        self.output_path.trace_add("write", lambda *_args: self._mark_output_stale())
        self._update_confidence_label()

    def _set_app_icon(self) -> None:
        if APP_ICON_ICO.exists():
            try:
                self.iconbitmap(default=str(APP_ICON_ICO))
            except tk.TclError:
                pass
        if APP_ICON_PNG.exists():
            try:
                self._app_icon_photo = tk.PhotoImage(file=str(APP_ICON_PNG))
                self.iconphoto(True, self._app_icon_photo)
            except tk.TclError:
                self._app_icon_photo = None

    def _build_ui(self) -> None:
        self.columnconfigure(0, weight=1)
        self.rowconfigure(0, weight=1)

        root = ttk.Frame(self, padding=16)
        root.pack(fill=tk.BOTH, expand=True)

        form = ttk.Frame(root)
        form.pack(fill=tk.X)
        form.columnconfigure(1, weight=1)

        ttk.Label(form, text="Bild").grid(row=0, column=0, sticky="w", padx=(0, 8), pady=4)
        ttk.Entry(form, textvariable=self.input_path).grid(row=0, column=1, sticky="ew", pady=4)
        ttk.Button(form, text="Auswaehlen...", command=self._pick_input).grid(row=0, column=2, padx=(8, 0), pady=4)

        ttk.Label(form, text="PowerPoint").grid(row=1, column=0, sticky="w", padx=(0, 8), pady=4)
        ttk.Entry(form, textvariable=self.output_path).grid(row=1, column=1, sticky="ew", pady=4)
        ttk.Button(form, text="Speichern als...", command=self._pick_output).grid(row=1, column=2, padx=(8, 0), pady=4)

        options = ttk.Frame(root)
        options.pack(fill=tk.X, pady=(12, 8))
        ttk.Label(options, text="OCR-Mindestkonfidenz").pack(side=tk.LEFT)
        ttk.Scale(
            options,
            from_=0.2,
            to=0.8,
            variable=self.confidence,
            orient=tk.HORIZONTAL,
            length=220,
            command=lambda _value: self._update_confidence_label(),
        ).pack(side=tk.LEFT, padx=10)
        ttk.Label(options, textvariable=self.confidence_label, width=4).pack(side=tk.LEFT)
        self.visuals_check = ttk.Checkbutton(options, text="Bildobjekte erkennen", variable=self.include_visuals)
        self.visuals_check.pack(side=tk.LEFT, padx=16)
        ttk.Label(options, text="Modus").pack(side=tk.LEFT, padx=(0, 6))
        self.mode_select = ttk.Combobox(
            options,
            textvariable=self.export_mode,
            values=EXPORT_MODES,
            width=18,
            state="readonly",
        )
        self.mode_select.pack(side=tk.LEFT)
        self.convert_button = ttk.Button(options, text="PPTX erzeugen", command=self._start_conversion)
        self.convert_button.pack(side=tk.RIGHT)
        self.open_output_button = ttk.Button(
            options,
            text="Oeffnen",
            command=self._open_output,
            state=tk.DISABLED,
        )
        self.open_output_button.pack(side=tk.RIGHT, padx=(0, 8))

        preview_frame = ttk.LabelFrame(root, text="Vorschau")
        preview_frame.pack(fill=tk.BOTH, expand=True, pady=(8, 8))
        self.preview = ttk.Label(preview_frame, anchor=tk.CENTER, text="Noch kein Bild ausgewaehlt")
        self.preview.pack(fill=tk.BOTH, expand=True, padx=8, pady=8)
        self.preview.bind("<Configure>", self._resize_preview)

        log_frame = ttk.LabelFrame(root, text="Status")
        log_frame.pack(fill=tk.BOTH, expand=False)
        self.log = tk.Text(log_frame, height=8, wrap=tk.WORD)
        self.log.pack(fill=tk.BOTH, expand=True, padx=8, pady=8)
        self.log.configure(state=tk.DISABLED)

        bottom = ttk.Frame(root)
        bottom.pack(fill=tk.X, pady=(8, 0))
        self.progress = ttk.Progressbar(bottom, mode="indeterminate")
        self.progress.pack(side=tk.LEFT, fill=tk.X, expand=True)
        ttk.Label(bottom, textvariable=self.status_text, width=32, anchor=tk.E).pack(side=tk.RIGHT, padx=(12, 0))

    def _bind_shortcuts(self) -> None:
        self.bind("<Control-o>", lambda _event: self._pick_input())
        self.bind("<Control-s>", lambda _event: self._pick_output())
        self.bind("<Return>", lambda _event: self._start_conversion())

    def _update_confidence_label(self) -> None:
        self.confidence_label.set(f"{self.confidence.get():.2f}")

    def _pick_input(self) -> None:
        path = filedialog.askopenfilename(
            title="Bild auswaehlen",
            filetypes=[("Bilder", "*.png *.jpg *.jpeg *.bmp *.tif *.tiff"), ("Alle Dateien", "*.*")],
        )
        if not path:
            return
        self.input_path.set(path)
        if not self.output_path.get():
            self.output_path.set(str(Path(path).with_suffix(".pptx")))
        self._load_preview(Path(path))

    def _pick_output(self) -> None:
        path = filedialog.asksaveasfilename(
            title="PowerPoint speichern",
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
        )
        if path:
            self.output_path.set(path)

    def _load_preview(self, path: Path) -> None:
        try:
            self._preview_source = Image.open(path).convert("RGB")
            self._resize_preview()
            self.status_text.set(f"Bild geladen: {self._preview_source.width} x {self._preview_source.height} px")
        except Exception as exc:
            self._preview_source = None
            self.preview.configure(image="", text="Vorschau nicht verfuegbar")
            messagebox.showerror("Vorschau fehlgeschlagen", str(exc))

    def _resize_preview(self, _event: tk.Event | None = None) -> None:
        if self._preview_source is None:
            return
        width = max(160, self.preview.winfo_width() - 16)
        height = max(120, self.preview.winfo_height() - 16)
        image = self._preview_source.copy()
        image.thumbnail((width, height))
        self._preview_image = ImageTk.PhotoImage(image)
        self.preview.configure(image=self._preview_image, text="")

    def _append_log(self, message: str) -> None:
        self.log.configure(state=tk.NORMAL)
        self.log.insert(tk.END, message + "\n")
        self.log.see(tk.END)
        self.log.configure(state=tk.DISABLED)
        self.status_text.set(message)

    def _thread_log(self, message: str) -> None:
        self.after(0, self._append_log, message)

    def _start_conversion(self) -> None:
        if self._is_running:
            return
        input_text = self.input_path.get().strip()
        output_text = self.output_path.get().strip()
        if not input_text:
            messagebox.showerror("Eingabe fehlt", "Bitte ein Bild auswaehlen.")
            return
        if not output_text:
            messagebox.showerror("Ausgabe fehlt", "Bitte einen PowerPoint-Ausgabepfad angeben.")
            return

        input_path = Path(input_text)
        output_path = Path(output_text)
        if not input_path.is_file():
            messagebox.showerror("Eingabe fehlt", "Bitte ein vorhandenes Bild auswaehlen.")
            return
        if output_path.suffix.lower() != ".pptx":
            output_path = output_path.with_suffix(".pptx")
            self.output_path.set(str(output_path))
        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            messagebox.showerror("Ausgabe nicht moeglich", f"Der Zielordner konnte nicht angelegt werden:\n{exc}")
            return

        self._set_running(True)
        self.open_output_button.configure(state=tk.DISABLED)
        self._append_log("Starte Umwandlung ...")
        worker = threading.Thread(
            target=self._run_conversion,
            args=(
                input_path,
                output_path,
                float(self.confidence.get()),
                bool(self.include_visuals.get()),
                self.export_mode.get(),
            ),
            daemon=True,
        )
        worker.start()

    def _run_conversion(
        self, input_path: Path, output_path: Path, confidence: float, include_visuals: bool, mode: str
    ) -> None:
        try:
            self._thread_log("Analysiere Folien fuer den Editor ...")
            slides = analyze_image(input_path, confidence, include_visuals, self._thread_log, DEFAULT_CONFIG, None)
        except Exception as exc:
            self.after(0, messagebox.showerror, "Analyse fehlgeschlagen", str(exc))
            self._thread_log(f"Fehler: {exc}")
            self.after(0, self._set_running, False)
            return
        self._pending_output_path = output_path
        self._pending_mode = mode
        self.after(0, self._open_reconstruction_editor, slides)

    def _open_reconstruction_editor(self, slides: list[AnalyzedSlide]) -> None:
        self._set_running(False)
        self._append_log("Editor geoeffnet. Inhalte pruefen und dann PPTX erzeugen.")
        ReconstructionEditor(self, slides, self._export_edited_slides)

    def _export_edited_slides(self, slides: list[AnalyzedSlide]) -> None:
        output_path = self._pending_output_path
        if output_path is None:
            messagebox.showerror("Ausgabe fehlt", "Kein PowerPoint-Ausgabepfad vorhanden.")
            return
        self._set_running(True)
        self.open_output_button.configure(state=tk.DISABLED)
        self._append_log("Erzeuge PowerPoint aus bearbeiteter Rekonstruktion ...")
        worker = threading.Thread(
            target=self._run_export,
            args=(slides, output_path, self._pending_mode),
            daemon=True,
        )
        worker.start()

    def _run_export(self, slides: list[AnalyzedSlide], output_path: Path, mode: str) -> None:
        try:
            export_pptx(slides, output_path, self._thread_log, mode, DEFAULT_CONFIG, None)
        except Exception as exc:
            self.after(0, messagebox.showerror, "Umwandlung fehlgeschlagen", str(exc))
            self._thread_log(f"Fehler: {exc}")
            self.after(0, self._set_running, False)
            return
        self._thread_log(f"Fertig: {output_path}")
        self.after(0, self._conversion_finished)
        self.after(0, messagebox.showinfo, "Fertig", f"PowerPoint wurde erzeugt:\n{output_path}")

    def _set_running(self, is_running: bool) -> None:
        self._is_running = is_running
        state = tk.DISABLED if is_running else tk.NORMAL
        self.convert_button.configure(state=state)
        self.visuals_check.configure(state=state)
        self.mode_select.configure(state=tk.DISABLED if is_running else "readonly")
        if is_running:
            self.progress.start(12)
        else:
            self.progress.stop()

    def _conversion_finished(self) -> None:
        self._set_running(False)
        if Path(self.output_path.get()).exists():
            self.open_output_button.configure(state=tk.NORMAL)

    def _mark_output_stale(self) -> None:
        if hasattr(self, "open_output_button"):
            self.open_output_button.configure(state=tk.DISABLED)

    def _open_output(self) -> None:
        path = Path(self.output_path.get())
        if not path.exists():
            messagebox.showerror("Datei fehlt", "Die PowerPoint-Datei wurde nicht gefunden.")
            return
        os.startfile(path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert a slide image grid to an editable PPTX.")
    parser.add_argument("--input", "-i", type=Path)
    parser.add_argument("--output", "-o", type=Path)
    parser.add_argument("--confidence", type=float, default=MIN_TEXT_CONFIDENCE)
    parser.add_argument("--no-visuals", action="store_true")
    parser.add_argument("--mode", choices=EXPORT_MODES, default="hybrid")
    parser.add_argument("--config", type=Path, help="Optionale JSON-Konfigurationsdatei.")
    parser.add_argument("--debug-dir", type=Path, help="Verzeichnis fuer Debug-Bilder, SVGs und Reports.")
    parser.add_argument("--debug", action="store_true", help="Debug-Ausgaben neben der PPTX erzeugen.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if args.input and args.output:
        config = load_config(args.config)
        if args.debug:
            config = ReconstructionConfig(**{**config.__dict__, "enable_debug_output": True})
        convert_image_to_pptx(
            args.input,
            args.output,
            confidence=args.confidence,
            include_visuals=not args.no_visuals,
            progress=print,
            mode=args.mode,
            config=config,
            debug_dir=args.debug_dir,
        )
        return

    configure_windows_taskbar_icon()
    app = Image2PptApp()
    show_startup_splash(app)
    app.mainloop()


if __name__ == "__main__":
    main()
